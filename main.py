import os, tempfile, time
from flask import Flask, request, render_template, send_file, jsonify, flash, redirect, url_for, Response
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
# from docx2pdf import convert
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from PIL import Image
import PyPDF2
import io
import traceback
from advanced_text_filter import filter_text_blocks
import re
from document_manager import DocumentManager

# OCR 관련 import 추가
try:
    import pytesseract
    import cv2
    import numpy as np
    OCR_AVAILABLE = True
    print("✅ OCR 모듈 로드 성공")
except ImportError as e:
    OCR_AVAILABLE = False
    print(f"⚠️ OCR 모듈 로드 실패: {e}")
    print("📥 pip install pytesseract opencv-python 실행 필요")

# Tesseract 경로 설정 (Windows)
if OCR_AVAILABLE:
    try:
        # Windows Tesseract 경로 자동 감지
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            r'C:\Users\{}\AppData\Local\Tesseract-OCR\tesseract.exe'.format(os.getenv('USERNAME', '')),
            'tesseract'  # PATH에 있는 경우
        ]
        
        for path in possible_paths:
            if os.path.exists(path) or path == 'tesseract':
                pytesseract.pytesseract.tesseract_cmd = path
                print(f"✅ Tesseract 경로 설정: {path}")
                break
        else:
            print("⚠️ Tesseract 실행 파일을 찾을 수 없습니다.")
            print("📥 https://github.com/UB-Mannheim/tesseract/wiki 에서 다운로드 필요")
    except Exception as e:
        print(f"⚠️ Tesseract 설정 오류: {e}")

app = Flask(__name__) 
CORS(app)  # CORS 활성화
app.secret_key = 'your-secret-key-here'
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB 제한 
UPLOAD_FOLDER = 'uploads' 
OUTPUT_FOLDER = 'outputs' 
os.makedirs(UPLOAD_FOLDER, exist_ok=True) 
os.makedirs(OUTPUT_FOLDER, exist_ok=True) 

# 전역 변수로 DocumentManager 인스턴스 생성
doc_manager = DocumentManager()

def clean_extracted_text(text):
    """추출된 텍스트를 정리하는 함수 (풋터 중복 및 메타데이터 완전 제거)"""
    if not text:
        return ""
    
    try:
        import re
        
        # 1단계: 기본 정리
        cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        cleaned = re.sub(r'\s+', ' ', cleaned)
        cleaned = cleaned.strip()
        
        if not cleaned:
            return ""
        
        # 2단계: 🔴 풋터 메타데이터 완전 제거 (강화된 패턴)
        footer_patterns = [
            # 기본 풋터 패턴들
            r"MARSCONTEST.*?의학전문대학원입학추가시험.*?부산대.*?메이커발명융합연구소.*?박상규.*?선임연구원.*?\(010-3577-1389\).*?붙임.*?참여.*?강사.*?별첨하오니.*?각.*?1부.*?붙임.*?행사.*?영문.*?2025.*?30일.*?수.*?47503.*?부산광역시.*?연제구.*?금련로.*?24.*?거제동.*?부산대학교.*?www\.bnue\.ac\.kr.*?전화번호.*?051-512-8803.*?팩스번호.*?051-512-0367.*?전자",
            
            # 세분화된 풋터 패턴들
            r"MARSCONTEST.*?의학전문대학원입학추가시험",
            r"부산대.*?메이커발명융합연구소",
            r"박상규.*?선임연구원.*?\(010-3577-1389\)",
            r"붙임.*?참여.*?강사.*?별첨하오니.*?각.*?1부",
            r"붙임.*?행사.*?영문.*?2025.*?30일.*?수",
            r"47503.*?부산광역시.*?연제구.*?금련로.*?24",
            r"거제동.*?부산대학교",
            r"www\.bnue\.ac\.kr",
            r"전화번호.*?051-512-8803.*?팩스번호.*?051-512-0367",
            r"전자.*?메일.*?www\.bnue\.ac\.kr",
            
            # 개별 요소들
            r"051-512-8803",
            r"051-512-0367",
            r"010-3577-1389",
            r"부산광역시.*?연제구.*?금련로.*?24",
            r"거제동.*?부산대학교",
            r"메이커발명융합연구소",
            r"의학전문대학원입학추가시험",
            r"선임연구원.*?박상규",
            r"박상규.*?선임연구원",
            
            # 일반적인 풋터 패턴들
            r"전화\s*[:：]?\s*\d{2,3}-\d{3,4}-\d{4}",
            r"팩스\s*[:：]?\s*\d{2,3}-\d{3,4}-\d{4}",
            r"휴대폰\s*[:：]?\s*\d{3}-\d{4}-\d{4}",
            r"이메일\s*[:：]?\s*[\w.-]+@[\w.-]+\.[a-zA-Z]{2,}",
            r"홈페이지\s*[:：]?\s*www\.[\w.-]+",
            r"주소\s*[:：]?\s*\d{5}.*?[시도].*?[구군].*?[동읍면]",
            r"우편번호\s*[:：]?\s*\d{5}",
            
            # 반복되는 연락처 정보
            r"담당자\s*[:：]?.*?\d{3}-\d{4}-\d{4}",
            r"연락처\s*[:：]?.*?\d{2,3}-\d{3,4}-\d{4}",
            r"문의\s*[:：]?.*?\d{2,3}-\d{3,4}-\d{4}",
            
            # 기관 정보
            r"[가-힣]+대학교.*?[가-힣]+연구소",
            r"[가-힣]+연구소.*?[가-힣]+대학교",
            r"[가-힣]+학과.*?[가-힣]+대학",
            
            # 날짜 및 시간 정보
            r"\d{4}년\s*\d{1,2}월\s*\d{1,2}일",
            r"\d{4}\.\s*\d{1,2}\.\s*\d{1,2}",
            r"\d{4}-\d{1,2}-\d{1,2}",
            
            # 페이지 번호 및 문서 정보
            r"페이지\s*\d+\s*/\s*\d+",
            r"\d+\s*/\s*\d+\s*페이지",
            r"- \d+ -",
            r"\[\s*\d+\s*\]",
            
            # 기타 메타데이터
            r"작성자\s*[:：]?.*?수정일\s*[:：]?",
            r"생성일\s*[:：]?.*?수정일\s*[:：]?",
            r"파일명\s*[:：]?.*?\.pdf",
            r"문서번호\s*[:：]?.*?-\d+"
        ]
        
        # 풋터 패턴들 제거
        for pattern in footer_patterns:
            cleaned = re.sub(pattern, '', cleaned, flags=re.IGNORECASE | re.DOTALL)
        
        # 3단계: 🔄 페이지 중복 제거 (라인 단위)
        lines = [line.strip() for line in cleaned.split('\n') if line.strip()]
        if not lines:
            return ""
        
        # 중복 라인 제거 (더 강화된 방식)
        unique_lines = []
        seen_lines = set()
        
        for line in lines:
            # 정규화된 라인으로 중복 체크 (공백, 특수문자, 숫자 제거)
            normalized_line = re.sub(r'[\s\W\d]+', '', line.lower())
            
            # 너무 짧은 라인이나 중복 라인 제거
            if len(normalized_line) > 3 and normalized_line not in seen_lines:
                # 추가 풋터 키워드 체크
                footer_keywords = [
                    'marscontest', '부산대', '메이커발명', '연구소', '박상규', 
                    '선임연구원', '전화번호', '팩스번호', '주소', '우편번호',
                    '연락처', '담당자', '문의', '이메일', '홈페이지'
                ]
                
                # 풋터 키워드가 포함된 라인 제거
                if not any(keyword in line.lower() for keyword in footer_keywords):
                    seen_lines.add(normalized_line)
                    unique_lines.append(line)
        
        # 4단계: 📝 텍스트 중복 제거 (문장 단위)
        if unique_lines:
            text_content = ' '.join(unique_lines)
            sentences = re.split(r'[.!?]\s+', text_content)
            unique_sentences = []
            seen_sentences = set()
            
            for sentence in sentences:
                sentence = sentence.strip()
                if len(sentence) < 10:  # 너무 짧은 문장 제거
                    continue
                
                # 정규화된 문장으로 중복 체크
                normalized = re.sub(r'[\s\W\d]+', '', sentence.lower())
                if len(normalized) > 5 and normalized not in seen_sentences:
                    # 풋터 관련 문장 제거
                    footer_sentence_keywords = [
                        'marscontest', '부산대', '메이커발명', '연구소', '박상규',
                        '전화', '팩스', '주소', '연락', '담당', '문의'
                    ]
                    
                    if not any(keyword in sentence.lower() for keyword in footer_sentence_keywords):
                        seen_sentences.add(normalized)
                        unique_sentences.append(sentence)
            
            # 최대 4개 문장만 보존 (더 짧게)
            unique_sentences = unique_sentences[:4]
        else:
            unique_sentences = []
        
        # 5단계: 📏 최종 조립 및 길이 제한
        if unique_sentences:
            result_text = '. '.join(unique_sentences)
            if result_text and not result_text.endswith(('.', '!', '?')):
                result_text += '.'
        else:
            # 문장 분리가 실패한 경우 라인 기반으로 처리 (더 짧게)
            result_text = ' '.join(unique_lines[:8])  # 최대 8줄
        
        # 길이 제한 (600자로 더 짧게)
        max_length = 600
        if len(result_text) > max_length:
            # 문장 경계에서 자르기
            truncated = result_text[:max_length]
            last_sentence_end = max(truncated.rfind('.'), truncated.rfind('!'), truncated.rfind('?'))
            if last_sentence_end > max_length * 0.5:  # 50% 이상 지점에서 문장이 끝나면
                result_text = truncated[:last_sentence_end + 1]
            else:
                result_text = truncated.rstrip() + "..."
        
        # 6단계: 🧹 최종 정리
        result_text = re.sub(r'\s+', ' ', result_text)  # 연속 공백 제거
        result_text = re.sub(r'\s*\.\s*\.\s*\.+', '...', result_text)  # 연속 점 정리
        result_text = result_text.strip()
        
        # 7단계: 최종 검증 (풋터 키워드 재확인)
        if result_text:
            final_footer_check = [
                'marscontest', '부산대', '메이커발명', '연구소', '박상규',
                '051-512', '010-3577', 'www.bnue', '부산광역시', '연제구'
            ]
            
            # 풋터 키워드가 여전히 많이 포함되어 있으면 더 간단하게 처리
            footer_count = sum(1 for keyword in final_footer_check 
                             if keyword in result_text.lower())
            
            if footer_count > 2:  # 풋터 키워드가 2개 이상이면
                # 공문서 핵심 키워드만 추출
                doc_keywords = ['공지', '안내', '수신', '제목', '담당', '회의', '검토', '요청', '신청']
                doc_sentences = []
                
                for sentence in unique_sentences:
                    if any(keyword in sentence for keyword in doc_keywords):
                        doc_sentences.append(sentence)
                
                if doc_sentences:
                    result_text = '. '.join(doc_sentences[:2])  # 최대 2개 문장
                    if not result_text.endswith('.'):
                        result_text += '.'
                else:
                    # 풋터가 너무 많으면 빈 텍스트 반환
                    return ""
        
        return result_text if len(result_text) > 15 else ""
        
    except Exception as e:
        print(f"텍스트 정리 중 오류: {e}")
        return ""

def allowed_file(filename): 
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pdf', 'docx'} 

def get_blank_slide_layout(prs):
    """안전한 빈 슬라이드 레이아웃 가져오기"""
    try:
        # 일반적으로 6번이 빈 레이아웃이지만, 없으면 다른 것 사용
        if len(prs.slide_layouts) > 6:
            return prs.slide_layouts[6]  # 빈 슬라이드
        elif len(prs.slide_layouts) > 5:
            return prs.slide_layouts[5]  # 제목만 있는 슬라이드
        elif len(prs.slide_layouts) > 0:
            return prs.slide_layouts[0]  # 첫 번째 사용 가능한 레이아웃
        else:
            # 레이아웃이 없으면 기본 생성
            return prs.slide_layouts[0]
    except IndexError:
        # 모든 경우에 실패하면 첫 번째 레이아웃 사용
        return prs.slide_layouts[0]

def pdf_to_pptx_with_images_only(pdf_path, output_path, quality='medium'):
    """이미지만 사용한 PDF → PPTX 변환 (텍스트 오버레이 없음)"""
    try:
        print(f"🖼️ 이미지 전용 PDF → PPTX 변환 시작: {pdf_path}")
        
        # 품질에 따른 DPI 설정
        dpi_settings = {
            'high': 300,
            'medium': 200,
            'low': 150
        }
        dpi = dpi_settings.get(quality, 200)
        
        # PDF를 이미지로 변환
        images = convert_from_path(pdf_path, dpi=dpi)
        print(f"📸 {len(images)}개 페이지 이미지 변환 완료 (DPI: {dpi})")
        
        # PPTX 생성
        prs = Presentation()
        
        # 첫 번째 이미지로 기본 슬라이드 크기 설정
        if images:
            first_image = images[0]
            if first_image.width > first_image.height:
                prs.slide_width = Inches(11.69)  # 가로형
                prs.slide_height = Inches(8.27)
                print("📐 가로형 슬라이드로 설정")
            else:
                prs.slide_width = Inches(8.27)   # 세로형
                prs.slide_height = Inches(11.69)
                print("📐 세로형 슬라이드로 설정")
        
        for i, image in enumerate(images):
            print(f"🔄 페이지 {i + 1}/{len(images)} 처리 중...")
            
            # 슬라이드 생성 - 안전한 레이아웃 사용
            slide_layout = get_blank_slide_layout(prs)  # 빈 슬라이드
            slide = prs.slides.add_slide(slide_layout)
            
            # 이미지를 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                # 슬라이드 크기에 맞게 이미지 크기 조정
                slide_width_inches = prs.slide_width.inches
                slide_height_inches = prs.slide_height.inches
                
                # 원본 이미지 크기 (인치)
                original_width_inches = image.width / dpi
                original_height_inches = image.height / dpi
                
                # 슬라이드에 맞게 스케일 조정
                scale_factor = min(
                    slide_width_inches / original_width_inches,
                    slide_height_inches / original_height_inches
                )
                
                final_width_inches = original_width_inches * scale_factor
                final_height_inches = original_height_inches * scale_factor
                
                # 이미지 리사이즈
                final_width_pixels = int(final_width_inches * dpi)
                final_height_pixels = int(final_height_inches * dpi)
                
                resized_image = image.resize(
                    (final_width_pixels, final_height_pixels), 
                    Image.Resampling.LANCZOS
                )
                
                # PNG로 저장
                resized_image.save(temp_file.name, 'PNG', optimize=True)
                
                # 슬라이드에 이미지 추가 (중앙 정렬)
                left = (prs.slide_width - Inches(final_width_inches)) / 2
                top = (prs.slide_height - Inches(final_height_inches)) / 2
                
                slide.shapes.add_picture(
                    temp_file.name, left, top, 
                    width=Inches(final_width_inches), 
                    height=Inches(final_height_inches)
                )
                
                print(f"✅ 페이지 {i + 1}: 이미지 추가 완료 (텍스트 오버레이 없음)")
                
                # 임시 파일 삭제
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
        
        # PPTX 저장
        prs.save(output_path)
        print(f"💾 이미지 전용 PPTX 저장 완료: {output_path}")
        
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"🎉 이미지 전용 변환 성공: {len(images)}개 페이지, {file_size:,} bytes")
            return True
        else:
            print(f"❌ 출력 파일이 생성되지 않았습니다")
            return False
        
    except Exception as e:
        print(f"❌ 이미지 전용 변환 중 오류: {e}")
        print(f"📍 오류 위치: {traceback.format_exc()}")
        return False

def extract_korean_optimized_blocks(ocr_data):
    """극한 완화된 텍스트 블록 추출"""
    blocks = []
    for i in range(len(ocr_data['text'])):
        if int(ocr_data['conf'][i]) < 5:  # 5%까지 극한 완화!
            continue
        text = ocr_data['text'][i].strip()
        if not text:  # 빈 텍스트만 제외
            continue
        
        # 매우 관대한 텍스트 체크
        has_korean = any('가' <= c <= '힣' or 'ㄱ' <= c <= 'ㅎ' or 'ㅏ' <= c <= 'ㅣ' for c in text)
        has_english = any(c.isalpha() for c in text)
        has_number = any(c.isdigit() for c in text)
        has_symbol = any(c in '.,()[]{}:;-_/\"\' ' for c in text)
        
        # 하나라도 해당하면 허용
        if has_korean or has_english or has_number or has_symbol:
            blocks.append({
                'text': text,
                'x': int(ocr_data['left'][i]),
                'y': int(ocr_data['top'][i]),
                'w': int(ocr_data['width'][i]),
                'h': int(ocr_data['height'][i]),
                'conf': int(ocr_data['conf'][i])
            })
            print(f"🔍 극한추출: '{text}' (신뢰도: {int(ocr_data['conf'][i])}%)")
    
    print(f"📊 극한 추출 결과: {len(blocks)}개 블록")
    return sorted(blocks, key=lambda x: (x['y'], x['x']))

def filter_ocr_blocks(ocr_data, confidence_threshold):
    """OCR 데이터에서 유효한 블록 필터링 (메모리 효율적)"""
    valid_blocks = []
    
    try:
        for i, conf in enumerate(ocr_data['conf']):
            if int(conf) >= confidence_threshold and ocr_data['text'][i].strip():
                text = ocr_data['text'][i].strip()
                if len(text) >= 1:
                    has_korean = any('가' <= c <= '힣' or 'ㄱ' <= c <= 'ㅎ' or 'ㅏ' <= c <= 'ㅣ' for c in text)
                    has_english = any(c.isalpha() for c in text)
                    has_number = any(c.isdigit() for c in text)
                    has_symbol = any(c in '.,()[]{}:;-_/\\"\' ' for c in text)
                    
                    if has_korean or has_english or has_number or has_symbol:
                        valid_blocks.append({
                            'text': text,
                            'x': int(ocr_data['left'][i]),
                            'y': int(ocr_data['top'][i]),
                            'w': int(ocr_data['width'][i]),
                            'h': int(ocr_data['height'][i]),
                            'conf': int(conf)
                        })
    except Exception as e:
        print(f"⚠️ 블록 필터링 오류: {e}")
    
    return valid_blocks

def merge_adjacent_blocks(blocks):
    """인접한 텍스트 블록 병합 (메모리 효율적)"""
    if not blocks or len(blocks) == 0:
        return blocks
    
    print(f"🔗 블록 병합 시작: {len(blocks)}개 → ", end="")
    
    try:
        # Y좌표 순으로 정렬
        sorted_blocks = sorted(blocks, key=lambda x: (x['y'], x['x']))
        merged_blocks = []
        
        i = 0
        while i < len(sorted_blocks):
            current_block = sorted_blocks[i]
            
            # 같은 줄에 있는 블록들 찾기 (Y좌표 차이 15px 이내로 축소)
            same_line_blocks = [current_block]
            j = i + 1
            
            while j < len(sorted_blocks) and j < i + 5:  # 최대 5개까지만 체크
                next_block = sorted_blocks[j]
                
                # 같은 줄 판정 (Y좌표 차이 15px 이내)
                if abs(next_block['y'] - current_block['y']) <= 15:
                    # X좌표 차이 80px 이내면 병합 대상
                    if abs(next_block['x'] - (current_block['x'] + current_block['w'])) <= 80:
                        same_line_blocks.append(next_block)
                        j += 1
                    else:
                        break
                else:
                    break
            
            # 블록 병합
            if len(same_line_blocks) > 1:
                # 여러 블록을 하나로 병합
                merged_text = ' '.join([block['text'] for block in same_line_blocks])
                merged_x = min([block['x'] for block in same_line_blocks])
                merged_y = min([block['y'] for block in same_line_blocks])
                merged_w = max([block['x'] + block['w'] for block in same_line_blocks]) - merged_x
                merged_h = max([block['h'] for block in same_line_blocks])
                merged_conf = max([block['conf'] for block in same_line_blocks])
                
                merged_blocks.append({
                    'text': merged_text,
                    'x': merged_x,
                    'y': merged_y,
                    'w': merged_w,
                    'h': merged_h,
                    'conf': merged_conf
                })
                
                i = j
            else:
                # 단일 블록
                merged_blocks.append(current_block)
                i += 1
        
        print(f"{len(merged_blocks)}개")
        return merged_blocks
        
    except Exception as e:
        print(f"병합 오류: {e}")
        return blocks  # 오류 시 원본 반환

def extract_text_blocks_with_position(ocr_data):
    """OCR 데이터에서 텍스트 블록과 위치 정보 추출 (한국어 최적화)"""
    text_blocks = []
    seen_texts = set()  # 중복 텍스트 추적
    seen_positions = []  # 중복 위치 추적
    
    try:
        n_boxes = len(ocr_data['text'])
        print(f"🔍 총 {n_boxes}개 OCR 박스 분석 중...")
        
        # 먼저 유효한 블록 필터링
        filtered_blocks = filter_ocr_blocks(ocr_data, 30)  # 30% 신뢰도 임계값
        
        # 인접 블록 병합
        merged_blocks = merge_adjacent_blocks(filtered_blocks)
        
        for block in merged_blocks:
            text = block['text']
            
            # 🔥 중복 텍스트 제거 (더 관대하게)
            normalized_text = ''.join(text.split()).lower()
            if len(normalized_text) > 0 and normalized_text in seen_texts:
                print(f"⚠️ 중복 텍스트 제거: '{text[:15]}...'")
                continue
            
            # 위치 정보 추출
            x = block['x']
            y = block['y']
            width = block['w']
            height = block['h']
            confidence = block['conf']
            
            # 🔥 중복 위치 제거 (더 관대하게 - 90% → 85%)
            is_duplicate_position = False
            for prev_pos in seen_positions:
                overlap_x = max(0, min(x + width, prev_pos['x'] + prev_pos['width']) - max(x, prev_pos['x']))
                overlap_y = max(0, min(y + height, prev_pos['y'] + prev_pos['height']) - max(y, prev_pos['y']))
                overlap_area = overlap_x * overlap_y
                
                current_area = width * height
                prev_area = prev_pos['width'] * prev_pos['height']
                
                if current_area > 0 and prev_area > 0:
                    overlap_ratio = overlap_area / min(current_area, prev_area)
                    if overlap_ratio > 0.85:  # 80% → 85%로 완화
                        print(f"⚠️ 중복 위치 제거: '{text[:15]}...' (겹침률: {overlap_ratio:.1%})")
                        is_duplicate_position = True
                        break
            
            if is_duplicate_position:
                continue
            
            # 🔥 한국어 문자 포함 여부 체크
            has_korean = any('가' <= c <= '힣' or 'ㄱ' <= c <= 'ㅎ' or 'ㅏ' <= c <= 'ㅣ' for c in text)
            has_english = any(c.isalpha() for c in text)
            has_number = any(c.isdigit() for c in text)
            
            # 텍스트 블록 정보 생성
            text_block = {
                'text': text,
                'x': x,
                'y': y,
                'width': width,
                'height': height,
                'confidence': confidence,
                'has_korean': has_korean,
                'has_english': has_english,
                'has_number': has_number
            }
            
            # 중복 추적에 추가
            if len(normalized_text) > 0:
                seen_texts.add(normalized_text)
            seen_positions.append({
                'x': x, 'y': y, 'width': width, 'height': height
            })
            
            text_blocks.append(text_block)
            print(f"✅ 텍스트 추가: '{text[:20]}...' (신뢰도: {confidence}%, 한글: {has_korean})")
    
    except Exception as e:
        print(f"⚠️ OCR 데이터 처리 오류: {e}")
        import traceback
        print(f"📍 오류 상세: {traceback.format_exc()}")
        return []
    
    # Y좌표 순으로 정렬
    text_blocks.sort(key=lambda x: (x['y'], x['x']))
    
    print(f"📝 최종 추출: {len(text_blocks)}개 텍스트 블록 (한국어 최적화)")
    return text_blocks

def create_korean_optimized_slide(slide, blocks, slide_w, slide_h, img_size):
    """메모리 최적화된 슬라이드 생성"""
    try:
        print(f"🎨 슬라이드 생성 시작: {len(blocks)}개 블록")
        
        if not blocks or len(blocks) == 0:
            print("⚠️ 블록이 없어 빈 슬라이드 생성")
            return
        
        # 블록이 이미 필터된 형태인지 확인
        if isinstance(blocks[0], dict) and 'text' in blocks[0]:
            # 이미 필터된 블록 형태
            text_blocks = blocks
        else:
            # 기존 OCR 데이터 형태면 변환
            text_blocks = extract_korean_optimized_blocks(blocks)
        
        # 🔥 블록 수 제한 (메모리 절약)
        if len(text_blocks) > 50:  # 최대 50개 블록만 처리
            text_blocks = sorted(text_blocks, key=lambda x: x['conf'], reverse=True)[:50]
            print(f"⚠️ 블록 수 제한: 상위 50개만 처리")
        
        # 🔥 블록 병합
        merged_blocks = merge_adjacent_blocks(text_blocks)
        
        # Y좌표 순으로 정렬
        merged_blocks.sort(key=lambda x: (x['y'], x['x']))
        
        # 🔥 제목 후보 찾기 (가장 상단의 큰 블록)
        title_candidate = None
        if merged_blocks:
            # 상위 20% 영역에서 가장 큰 블록
            top_20_percent = img_size[1] * 0.2
            top_blocks = [b for b in merged_blocks if b['y'] <= top_20_percent]
            
            if top_blocks:
                title_candidate = max(top_blocks, key=lambda x: x['w'] * x['h'])
        
        # 슬라이드에 텍스트 추가 (최대 30개까지만)
        processed_count = 0
        for i, block in enumerate(merged_blocks):
            if processed_count >= 30:  # 최대 30개 텍스트박스
                print(f"⚠️ 텍스트박스 수 제한: 30개까지만 처리")
                break
                
            try:
                # 위치 계산
                x_ratio = block['x'] / img_size[0]
                y_ratio = block['y'] / img_size[1]
                w_ratio = block['w'] / img_size[0]
                h_ratio = block['h'] / img_size[1]
                
                # 슬라이드 좌표로 변환
                x = slide_w * x_ratio
                y = slide_h * y_ratio
                w = max(slide_w * w_ratio, slide_w * 0.1)  # 🔥 최소 너비 10%
                h = max(slide_h * h_ratio, slide_h * 0.02)
                
                # 경계 체크
                if x + w > slide_w:
                    w = slide_w - x - slide_w * 0.01
                if y + h > slide_h:
                    h = slide_h - y - slide_h * 0.01
                
                if w > slide_w * 0.05 and h > slide_h * 0.01:  # 최소 크기 체크
                    # 텍스트박스 추가
                    textbox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
                    textbox.text_frame.text = block['text']
                    
                    # 🔥 폰트 크기 설정
                    for paragraph in textbox.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if block == title_candidate:
                                # 제목: 18-20pt
                                run.font.size = Pt(18)
                                run.font.bold = True
                            else:
                                # 일반 텍스트: 최소 12pt
                                font_size = max(12, min(14, int(h * 0.5)))
                                run.font.size = Pt(font_size)
                    
                    processed_count += 1
                        
            except Exception as e:
                print(f"⚠️ 블록 {i} 처리 오류: {e}")
                continue
        
        print(f"✅ 슬라이드 생성 완료: {processed_count}개 텍스트박스")
        
    except Exception as e:
        print(f"❌ 슬라이드 생성 오류: {e}")
        # 오류 시 기본 텍스트 추가
        try:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            textbox.text_frame.text = "OCR 처리 중 오류 발생\n메모리 부족 또는 파일 크기 문제"
        except:
            pass


def create_slide_with_ocr_text(slide, text_blocks, slide_width, slide_height, image_width, image_height):
    """OCR 텍스트를 PPTX 슬라이드에 정확한 위치로 배치"""
    try:
        # 좌표 변환 비율 계산
        width_ratio = slide_width / image_width
        height_ratio = slide_height / image_height
        
        # 텍스트 블록들을 Y 좌표 순으로 정렬 (위에서 아래로)
        sorted_blocks = sorted(text_blocks, key=lambda x: x['y'])
        
        for block in sorted_blocks:
            try:
                # 이미지 좌표를 슬라이드 좌표로 변환
                left = Inches(block['x'] * width_ratio / 72)  # 72 DPI 기준
                top = Inches(block['y'] * height_ratio / 72)
                width = Inches(block['width'] * width_ratio / 72)
                height = Inches(block['height'] * height_ratio / 72)
                
                # 최소/최대 크기 제한
                if width < Inches(0.5):
                    width = Inches(0.5)
                if height < Inches(0.2):
                    height = Inches(0.2)
                
                # 슬라이드 경계 내에 있는지 확인
                if left + width > slide_width:
                    width = slide_width - left - Inches(0.1)
                if top + height > slide_height:
                    height = slide_height - top - Inches(0.1)
                
                if width > Inches(0.1) and height > Inches(0.1):
                    # 텍스트 박스 추가
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    
                    # 텍스트 추가
                    p = text_frame.paragraphs[0]
                    p.text = block['text']
                    
                    # 폰트 설정
                    font = p.font
                    font.name = 'NanumGothic'  # 나눔고딕 폰트
                    
                    # 신뢰도에 따른 폰트 크기 조정
                    if block['confidence'] > 80:
                        font.size = Pt(12)
                    elif block['confidence'] > 60:
                        font.size = Pt(11)
                    else:
                        font.size = Pt(10)
                    
                    # 텍스트 정렬
                    p.alignment = PP_ALIGN.LEFT
                    
                    print(f"✅ OCR 텍스트 추가: '{block['text'][:20]}...' (신뢰도: {block['confidence']}%)")
            
            except Exception as block_error:
                print(f"⚠️ 텍스트 블록 처리 오류: {block_error}")
                continue
    
    except Exception as e:
        print(f"❌ OCR 텍스트 배치 오류: {e}")

def combine_image_and_text(slide, image, text_blocks):
    """이미지 배경과 텍스트 레이어 결합"""
    try:
        # 이미지를 임시 파일로 저장
        temp_image_path = os.path.join(tempfile.gettempdir(), f"temp_slide_{int(time.time())}.png")
        image.save(temp_image_path, 'PNG')
        
        # 슬라이드에 배경 이미지 추가
        slide_width = slide.slide_layout.slide_master.slide_width
        slide_height = slide.slide_layout.slide_master.slide_height
        
        # 이미지를 슬라이드 크기에 맞게 추가
        slide.shapes.add_picture(
            temp_image_path, 
            0, 0, 
            slide_width, 
            slide_height
        )
        
        # 임시 파일 삭제
        try:
            os.remove(temp_image_path)
        except:
            pass
        
        print(f"✅ 배경 이미지 추가 완료")
    
    except Exception as e:
        print(f"⚠️ 이미지-텍스트 결합 오류: {e}")

def advanced_preprocess_for_korean_ocr(image):
    """한국어 OCR용 이미지 전처리 (메모리 최적화 포함)"""
    try:
        print("🔧 한국어 OCR 이미지 전처리 시작...")
        
        # PIL to OpenCV
        opencv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
        
        original_height, original_width = gray.shape
        original_pixels = original_width * original_height
        
        print(f"📏 원본 크기: {original_width}x{original_height} ({original_pixels:,} 픽셀)")
        
        # 🔥 픽셀 수 제한 (15백만 픽셀로 더 엄격하게)
        MAX_PIXELS = 15_000_000  # 25M → 15M으로 축소
        
        if original_pixels > MAX_PIXELS:
            # 축소 비율 계산
            scale_ratio = (MAX_PIXELS / original_pixels) ** 0.5
            new_width = int(original_width * scale_ratio)
            new_height = int(original_height * scale_ratio)
            
            gray = cv2.resize(gray, (new_width, new_height), interpolation=cv2.INTER_AREA)
            print(f"📉 파일 크기 축소 적용: {original_width}x{original_height} → {new_width}x{new_height}")
            print(f"   축소 비율: {scale_ratio:.3f}, 픽셀 수: {new_width*new_height:,}")
        
        # 🔥 조건부 확대 (짧은 변이 1200px 미만일 때만) - 더 엄격하게
        current_height, current_width = gray.shape
        min_dimension = min(current_width, current_height)
        
        # 확대 후 픽셀 수 체크
        if min_dimension < 1200:  # 1500 → 1200으로 축소
            enlarged_pixels = (current_width * 2) * (current_height * 2)
            if enlarged_pixels <= MAX_PIXELS:  # 확대 후에도 제한 내에 있으면
                enlarged = cv2.resize(gray, (current_width*2, current_height*2), interpolation=cv2.INTER_LANCZOS4)
                print(f"📈 2배 확대 적용: {current_width}x{current_height} → {current_width*2}x{current_height*2}")
                print(f"   확대 이유: 짧은 변 {min_dimension}px < 1200px")
                gray = enlarged
            else:
                print(f"⚠️ 확대 생략: 확대 시 픽셀 수 초과 ({enlarged_pixels:,} > {MAX_PIXELS:,})")
        else:
            print(f"⏭️ 확대 생략: 짧은 변 {min_dimension}px >= 1200px")
        
        final_height, final_width = gray.shape
        final_pixels = final_width * final_height
        print(f"🎯 최종 크기: {final_width}x{final_height} ({final_pixels:,} 픽셀)")
        
        # 🔥 메모리 효율적인 전처리
        try:
            # 노이즈 제거 (더 가벼운 필터)
            denoised = cv2.medianBlur(gray, 3)  # bilateralFilter 대신 medianBlur
            
            # 대비 향상 (더 작은 타일)
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(4,4))  # 8x8 → 4x4
            enhanced = clahe.apply(denoised)
            
            # 적응형 이진화 (더 작은 커널)
            binary = cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 10)  # 21,15 → 11,10
            
            # 모폴로지 연산 (더 작은 커널)
            kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (1, 1))  # 2x2 → 1x1
            cleaned = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)
            
            print("✅ 메모리 최적화 전처리 완료")
            return Image.fromarray(cleaned)
            
        except Exception as process_error:
            print(f"⚠️ 전처리 중 오류, 기본 이진화 적용: {process_error}")
            # 최소한의 처리만 수행
            binary = cv2.threshold(gray, 127, 255, cv2.THRESH_BINARY)[1]
            return Image.fromarray(binary)
        
    except Exception as e:
        print(f"❌ 이미지 전처리 오류: {e}")
        # 실패 시 원본을 작게 축소해서 반환
        try:
            small_image = image.resize((800, 600), Image.Resampling.LANCZOS)
            print("🔄 원본을 800x600으로 축소하여 반환")
            return small_image
        except:
            return image

def preprocess_image_for_ocr(image):
    """한국어 문서를 위한 고급 이미지 전처리"""
    try:
        print("🔧 한국어 문서용 이미지 전처리 시작...")
        
        # PIL Image를 OpenCV 형식으로 변환
        opencv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        
        # 1. 그레이스케일 변환
        gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
        print(f"📏 이미지 크기: {gray.shape[1]}x{gray.shape[0]}")
        
        # 2. 🔥 한국어 문서용 노이즈 제거 (더 부드럽게)
        denoised = cv2.bilateralFilter(gray, 5, 50, 50)  # 더 부드러운 필터
        
        # 3. 🔥 대비 향상 (한국어 문서에 최적화)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))  # 더 부드러운 대비
        enhanced = clahe.apply(denoised)
        
        # 4. 🔥 적응형 이진화 (한국어에 최적화)
        # 여러 방법 시도해서 최적 결과 선택
        binary_methods = [
            ('ADAPTIVE_THRESH_MEAN_C', cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY, 15, 10)),
            ('ADAPTIVE_THRESH_GAUSSIAN_C', cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 15, 10)),
            ('OTSU', cv2.threshold(enhanced, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1])
        ]
        
        # 가장 많은 텍스트 영역을 가진 방법 선택
        best_binary = None
        best_score = 0
        best_method = ""
        
        for method_name, binary_img in binary_methods:
            # 텍스트 영역 추정 (흰색 픽셀 비율)
            white_ratio = np.sum(binary_img == 255) / binary_img.size
            # 적절한 텍스트 비율 (10-70%)
            if 0.1 <= white_ratio <= 0.7:
                score = min(white_ratio, 1 - white_ratio)  # 균형잡힌 비율 선호
                if score > best_score:
                    best_score = score
                    best_binary = binary_img
                    best_method = method_name
        
        if best_binary is None:
            # 모든 방법이 실패하면 기본 방법 사용
            best_binary = cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
            best_method = "DEFAULT"
        
        print(f"✅ 최적 이진화 방법: {best_method} (점수: {best_score:.3f})")
        
        # 5. 🔥 모폴로지 연산 (한국어 문자 구조 개선)
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 1))
        cleaned = cv2.morphologyEx(best_binary, cv2.MORPH_CLOSE, kernel)
        
        # 6. 🔥 최종 노이즈 제거 (작은 점들 제거)
        kernel = np.ones((2,2), np.uint8)
        final_cleaned = cv2.morphologyEx(cleaned, cv2.MORPH_OPEN, kernel)
        
        # OpenCV 이미지를 PIL 이미지로 변환
        processed_image = Image.fromarray(final_cleaned)
        
        print("✅ 한국어 문서용 이미지 전처리 완료")
        return processed_image
    
    except Exception as e:
        print(f"⚠️ 이미지 전처리 오류: {e}")
        return image  # 전처리 실패 시 원본 반환

def pdf_to_pptx_with_ocr(pdf_path, output_path, quality='medium'):
    """한국어 최적화 OCR 기능이 포함된 PDF → PPTX 변환"""
    try:
        print(f"🔍 한국어 최적화 OCR 변환 시작: {pdf_path}")
        if not OCR_AVAILABLE:
            print("⚠️ OCR 모듈이 없어 변환 실패")
            return False
        
        # 혁신적 DPI 설정 (극한 해상도)
        dpi = {'low': 600, 'medium': 900, 'high': 1200}.get(quality, 900)  # 혁신적 해상도
        print(f"🔥 혁신적 DPI 설정: {dpi}")
        try:
            images = convert_from_path(pdf_path, dpi=dpi, fmt='PNG')
        except:
            images = convert_from_path(pdf_path, dpi=200, fmt='PNG')
        
        prs = Presentation()
        total_blocks, success_pages = 0, 0
        
        for page_num, image in enumerate(images):
            try:
                print(f"🔄 페이지 {page_num + 1}/{len(images)} 처리 중...")
                processed_img = advanced_preprocess_for_korean_ocr(image)
                
                # 표 문서 특화 OCR 설정 시도
                # 극한 OCR 설정 시도
                configs = [
                    {'name': '극한_표문서_전용', 'config': '--oem 3 --psm 6 -c preserve_interword_spaces=1 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyzㄱ-ㅎㅏ-ㅣ가-힣.,()[]{}:;-_/\"\' ', 'lang': 'kor+eng'},
                    {'name': '극한_전체분석', 'config': '--oem 3 --psm 3 -c preserve_interword_spaces=1', 'lang': 'kor+eng'},
                    {'name': '극한_자동분할', 'config': '--oem 3 --psm 4 -c preserve_interword_spaces=1', 'lang': 'kor+eng'},
                    {'name': '극한_단일블록', 'config': '--oem 3 --psm 8 -c preserve_interword_spaces=1', 'lang': 'kor+eng'},
                    {'name': '극한_세로텍스트', 'config': '--oem 3 --psm 12 -c preserve_interword_spaces=1', 'lang': 'kor+eng'},
                    {'name': '극한_단일단어', 'config': '--oem 3 --psm 7 -c preserve_interword_spaces=1', 'lang': 'kor+eng'},
                    {'name': '극한_원시분석', 'config': '--oem 3 --psm 13 -c preserve_interword_spaces=1', 'lang': 'kor+eng'},
                    {'name': '극한_한국어만', 'config': '--oem 3 --psm 6 -c preserve_interword_spaces=1', 'lang': 'kor'},
                    {'name': '극한_영어만', 'config': '--oem 3 --psm 6 -c preserve_interword_spaces=1', 'lang': 'eng'},
                    # 추가: 매우 관대한 설정
                    {'name': '극한_관대모드', 'config': '--oem 1 --psm 6 -c preserve_interword_spaces=1', 'lang': 'kor+eng'}
                ]
                
                best_data, best_count = None, 0
                for cfg in configs:
                    try:
                        ocr_data = pytesseract.image_to_data(processed_img, lang=cfg['lang'], 
                                                            output_type=pytesseract.Output.DICT, config=cfg['config'])
                        # 🔥 신뢰도 10% 이상, 길이 제한 완전 제거
                        valid_texts = []
                        for i, conf in enumerate(ocr_data['conf']):
                            if int(conf) >= 10 and ocr_data['text'][i].strip():  # 20% → 10%로 극한 완화!
                                text = ocr_data['text'][i].strip()
                                # 한 글자라도 의미있으면 허용
                                if len(text) >= 1:  # 길이 제한 완전 제거
                                    has_korean = any('가' <= c <= '힣' or 'ㄱ' <= c <= 'ㅎ' or 'ㅏ' <= c <= 'ㅣ' for c in text)
                                    has_english = any(c.isalpha() for c in text)
                                    has_number = any(c.isdigit() for c in text)
                                    has_symbol = any(c in '.,()[]{}:;-_/\"\' ' for c in text)
                                    if has_korean or has_english or has_number or has_symbol:
                                        valid_texts.append(text)
                        
                        if len(valid_texts) > best_count:
                            best_data, best_count = ocr_data, len(valid_texts)
                            print(f"  ✅ {cfg['name']}: {len(valid_texts)}개 텍스트")
                    except:
                        continue
                
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                if best_data and best_count > 0:
                    # 필터된 블록으로 슬라이드 생성
                    create_korean_optimized_slide(slide, best_data, prs.slide_width, prs.slide_height, image.size)
                    total_blocks += len(best_data)
                    success_pages += 1
                    print(f"✅ 페이지 {page_num + 1}: {len(best_data)}개 블록 → 슬라이드 생성 완료")
                else:
                    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                    textbox.text_frame.text = f"페이지 {page_num + 1}\n\nOCR 인식 실패\n\n• 이미지 품질 확인 필요\n• 스캔 해상도 향상 권장\n• 채택 설정: {best_config_name}"
                    print(f"❌ 페이지 {page_num + 1}: OCR 실패 (채택 설정: {best_config_name})")
                    
            except Exception as e:
                print(f"❌ 페이지 {page_num + 1} 오류: {e}")
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
                textbox.text_frame.text = f"페이지 {page_num + 1} 처리 오류"
        
        prs.save(output_path)
        success_rate = (success_pages / len(images)) * 100 if images else 0
        print(f"\n🎉 변환 완료! 성공률: {success_rate:.1f}% ({success_pages}/{len(images)}페이지)")
        print(f"📊 총 {total_blocks}개 텍스트 블록 추출")
        return True
        
    except Exception as e:
        print(f"❌ OCR 변환 오류: {e}")
        import traceback
        print(f"📍 오류 상세: {traceback.format_exc()}")
        return False

def pdf_to_pptx(pdf_path, output_path, quality='medium'):
    """PDF를 PPTX로 변환 (중복 방지 강화)"""
    try:
        print(f"🔄 PDF → PPTX 변환 시작: {pdf_path}")
        
        # 출력 파일이 이미 존재하면 삭제
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
                print(f"🗑️ 기존 파일 삭제: {output_path}")
            except Exception as e:
                print(f"⚠️ 기존 파일 삭제 실패: {e}")
        
        # 🎯 변환 방식 우선순위 (중복 방지)
        success = False
        conversion_method = ""
        
        # 1순위: 텍스트 기반 변환
        print("📝 텍스트 기반 변환 시도")
        success = pdf_to_pptx_with_text(pdf_path, output_path, quality)
        conversion_method = "텍스트 기반"
        
        if success:
            print(f"✅ 텍스트 변환 성공")
            return True
        else:
            print(f"⚠️ 텍스트 변환 실패 - OCR 변환으로 폴백")
        
        # 2순위: OCR 기능이 활성화된 경우 OCR 사용
        if OCR_AVAILABLE and not success:
            print("🔍 OCR 기능 활성화됨 - OCR 변환 시도")
            success = pdf_to_pptx_with_ocr(pdf_path, output_path, quality)
            conversion_method = "OCR 기반"
            
            if success:
                print(f"✅ OCR 변환 성공")
                return True
            else:
                print(f"⚠️ OCR 변환 실패 - 이미지 기반 변환으로 폴백")
        
        # 3순위: 이미지 기반 변환 (최종 폴백)
        if not success:
            print("🖼️ 이미지 기반 변환 시도 (최종 폴백)")
            success = pdf_to_pptx_with_images_only(pdf_path, output_path, quality)
            conversion_method = "이미지 기반"
        
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"✅ {conversion_method} 변환 성공: {output_path} (크기: {file_size:,} bytes)")
            return True
        else:
            print(f"❌ 모든 변환 방식 실패")
            return False
            
    except Exception as e:
        print(f"❌ PDF to PPTX 변환 오류: {e}")
        return False

def pdf_to_pptx_with_images_only(pdf_path, output_path, quality='medium'):
    """이미지만 사용한 PDF → PPTX 변환 (텍스트 오버레이 없음)"""
    try:
        print(f"🖼️ 이미지 전용 PDF → PPTX 변환 시작: {pdf_path}")
        
        # 품질에 따른 DPI 설정
        dpi_settings = {
            'high': 300,
            'medium': 200,
            'low': 150
        }
        dpi = dpi_settings.get(quality, 200)
        
        # PDF를 이미지로 변환
        images = convert_from_path(pdf_path, dpi=dpi)
        print(f"📸 {len(images)}개 페이지 이미지 변환 완료 (DPI: {dpi})")
        
        # PPTX 생성
        prs = Presentation()
        
        # 첫 번째 이미지로 기본 슬라이드 크기 설정
        if images:
            first_image = images[0]
            if first_image.width > first_image.height:
                prs.slide_width = Inches(11.69)  # 가로형
                prs.slide_height = Inches(8.27)
                print("📐 가로형 슬라이드로 설정")
            else:
                prs.slide_width = Inches(8.27)   # 세로형
                prs.slide_height = Inches(11.69)
                print("📐 세로형 슬라이드로 설정")
        
        for i, image in enumerate(images):
            print(f"🔄 페이지 {i + 1}/{len(images)} 처리 중...")
            
            # 슬라이드 생성
            slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
            slide = prs.slides.add_slide(slide_layout)
            
            # 이미지를 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                # 슬라이드 크기에 맞게 이미지 크기 조정
                slide_width_inches = prs.slide_width.inches
                slide_height_inches = prs.slide_height.inches
                
                # 원본 이미지 크기 (인치)
                original_width_inches = image.width / dpi
                original_height_inches = image.height / dpi
                
                # 슬라이드에 맞게 스케일 조정
                scale_factor = min(
                    slide_width_inches / original_width_inches,
                    slide_height_inches / original_height_inches
                )
                
                final_width_inches = original_width_inches * scale_factor
                final_height_inches = original_height_inches * scale_factor
                
                # 이미지 리사이즈
                final_width_pixels = int(final_width_inches * dpi)
                final_height_pixels = int(final_height_inches * dpi)
                
                resized_image = image.resize(
                    (final_width_pixels, final_height_pixels), 
                    Image.Resampling.LANCZOS
                )
                
                # PNG로 저장
                resized_image.save(temp_file.name, 'PNG', optimize=True)
                
                # 슬라이드에 이미지 추가 (중앙 정렬)
                left = (prs.slide_width - Inches(final_width_inches)) / 2
                top = (prs.slide_height - Inches(final_height_inches)) / 2
                
                slide.shapes.add_picture(
                    temp_file.name, left, top, 
                    width=Inches(final_width_inches), 
                    height=Inches(final_height_inches)
                )
                
                print(f"✅ 페이지 {i + 1}: 이미지 추가 완료 (텍스트 오버레이 없음)")
                
                # 임시 파일 삭제
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
        
        # PPTX 저장
        prs.save(output_path)
        print(f"💾 이미지 전용 PPTX 저장 완료: {output_path}")
        
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"🎉 이미지 전용 변환 성공: {len(images)}개 페이지, {file_size:,} bytes")
            return True
        else:
            print(f"❌ 출력 파일이 생성되지 않았습니다")
            return False
        
    except Exception as e:
        print(f"❌ 이미지 전용 변환 중 오류: {e}")
        print(f"📍 오류 위치: {traceback.format_exc()}")
        return False

def pdf_to_pptx_with_layout(pdf_path, output_path, quality='medium'):
    """pdfplumber를 사용한 정교한 레이아웃 보존 변환"""
    try:
        print("변환 설정: 정교한 레이아웃 보존 변환 (pdfplumber 사용)")
        print("PDF에서 테이블 및 레이아웃 정보 추출 중...")
        
        # 새 PowerPoint 프레젠테이션 생성
        prs = Presentation()
        
        with pdfplumber.open(pdf_path) as pdf:
            print(f"총 {len(pdf.pages)}페이지 처리 중...")
            
            for page_num, page in enumerate(pdf.pages):
                print(f"페이지 {page_num+1}/{len(pdf.pages)} 처리 중...")
                
                # 빈 슬라이드 추가
                slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
                slide = prs.slides.add_slide(slide_layout)
                
                # 1. 테이블 추출 시도 (안정적인 기본 설정)
                try:
                    # 기본 테이블 감지 (가장 안정적)
                    tables = page.find_tables()
                    
                    # 테이블이 감지되지 않으면 선 기반 감지 시도
                    if not tables:
                        tables = page.find_tables({
                            "vertical_strategy": "lines",
                            "horizontal_strategy": "lines"
                        })
                    
                    # 여전히 감지되지 않으면 텍스트 기반 감지 시도
                    if not tables:
                        tables = page.find_tables({
                            "vertical_strategy": "text",
                            "horizontal_strategy": "text"
                        })
                        
                except Exception as table_error:
                    print(f"테이블 감지 중 오류: {table_error}")
                    tables = []
                
                if tables:
                    print(f"페이지 {page_num+1}: {len(tables)}개의 테이블 발견")
                    
                    for table_idx, table in enumerate(tables):
                        try:
                            # 테이블을 PPTX 테이블로 변환
                            create_pptx_table(slide, table, page_num, table_idx)
                        except Exception as e:
                            print(f"테이블 {table_idx+1} 처리 중 오류: {e}")
                            continue
                else:
                    print(f"페이지 {page_num+1}: 테이블 없음, 텍스트 레이아웃 분석")
                    
                    # 2. 테이블이 없는 경우 고급 텍스트 레이아웃 분석
                    create_advanced_text_layout(slide, page, page_num)
        
        # PPTX 파일 저장
        prs.save(output_path)
        print("정교한 레이아웃 보존 변환 완료")
        return True
        
    except Exception as e:
        print(f"레이아웃 보존 변환 중 오류 발생: {str(e)}")
        return False

# 파일 유틸리티 함수

def is_file_locked(filepath):
    """파일이 다른 프로세스에 의해 잠겨있는지 확인"""
    if not os.path.exists(filepath):
        return False
    try:
        with open(filepath, 'a') as f:
            pass
        return False
    except IOError:
        return True

def generate_safe_filename(filename, extracted_numbers=None):
    """안전한 파일명 생성 + 번호 기반 명명 규칙"""
    import re
    from datetime import datetime
    
    # 기본 파일명 정리
    base_name = os.path.splitext(filename)[0]
    extension = os.path.splitext(filename)[1]
    
    # 🔥 번호 기반 파일명 생성
    if extracted_numbers:
        name_parts = []
        
        # 우선순위: KC번호 > 등록번호 > 문서번호
        if 'kc_number' in extracted_numbers:
            name_parts.append(f"KC_{extracted_numbers['kc_number']}")
        elif 'registration_number' in extracted_numbers:
            name_parts.append(f"REG_{extracted_numbers['registration_number']}")
        elif 'document_number' in extracted_numbers:
            name_parts.append(f"DOC_{extracted_numbers['document_number']}")
        
        # 날짜 추가
        if 'date' in extracted_numbers:
            date_clean = re.sub(r'[.-]', '', extracted_numbers['date'])
            name_parts.append(date_clean)
        else:
            # 현재 날짜 추가
            name_parts.append(datetime.now().strftime('%Y%m%d'))
        
        if name_parts:
            structured_name = '_'.join(name_parts)
            print(f"📝 구조화된 파일명: {structured_name}{extension}")
            return structured_name + extension
    
    # 폴백: 기본 안전 파일명
    safe_name = re.sub(r'[^\w\s-]', '', base_name)
    safe_name = re.sub(r'[-\s]+', '-', safe_name)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    
    return f"{safe_name}_{timestamp}{extension}"

def cleanup_temp_files(directory, pattern, max_age_hours=1):
    """오래된 임시 파일 정리"""
    import glob
    import time
    current_time = time.time()
    for filepath in glob.glob(os.path.join(directory, pattern)):
        try:
            file_age = current_time - os.path.getmtime(filepath)
            if file_age > max_age_hours * 3600:  # 시간 -> 초 변환
                os.remove(filepath)
                print(f"임시 파일 정리: {filepath}")
        except Exception as e:
            print(f"임시 파일 정리 실패: {filepath} - {e}")

def safe_file_operation(output_path, operation_func):
    """안전한 파일 작업 수행 (수정)"""
    try:
        # 작업 수행
        result = operation_func()
        return True, result
        
    except PermissionError as e:
        return False, f"권한 오류: {e}"
    except Exception as e:
        return False, f"작업 오류: {e}"

# PDF to PPTX 변환 함수 개선
def pdf_to_docx_safe(pdf_path, output_path, quality='medium'):
    """안전한 PDF → DOCX 변환"""
    try:
        print(f"PDF 변환 시작: {pdf_path} -> {output_path}")
        
        # 1. 출력 파일 잠금 확인
        if os.path.exists(output_path) and is_file_locked(output_path):
            print(f"⚠️ 출력 파일이 사용 중입니다: {output_path}")
            return False, "파일이 사용 중입니다. Word를 닫고 다시 시도하세요."
        
        # 2. 임시 파일명 생성
        temp_output = generate_safe_filename(f"temp_{os.path.basename(output_path)}")
        temp_output_path = os.path.join(os.path.dirname(output_path), temp_output)
        
        print(f"임시 파일 경로: {temp_output_path}")
        
        # 3. 변환 작업 수행
        try:
            success = pdf_to_docx(pdf_path, temp_output_path, quality)
            if not success:
                return False, "PDF 변환 중 오류가 발생했습니다."
        except Exception as e:
            print(f"변환 오류: {e}")
            return False, f"변환 중 오류: {str(e)}"
        
        # 4. 임시 파일 확인
        if not os.path.exists(temp_output_path):
            return False, "변환된 파일이 생성되지 않았습니다."
        
        # 5. 임시 파일을 최종 파일로 이동
        try:
            if os.path.exists(output_path):
                os.remove(output_path)
            
            os.rename(temp_output_path, output_path)
            print(f"✅ 파일 이동 완료: {output_path}")
            
            return True, "변환이 성공적으로 완료되었습니다."
            
        except Exception as e:
            print(f"파일 이동 오류: {e}")
            return False, f"파일 이동 중 오류: {str(e)}"
            
    except Exception as e:
        print(f"❌ PDF → DOCX 변환 중 치명적 오류: {e}")
        return False, f"변환 중 오류가 발생했습니다: {str(e)}"

def docx_to_pdf_safe(docx_path, output_path):
    """안전한 DOCX → PDF 변환"""
    try:
        print(f"DOCX 변환 시작: {docx_path} -> {output_path}")
        
        # 1. 출력 파일 잠금 확인
        if os.path.exists(output_path) and is_file_locked(output_path):
            print(f"⚠️ 출력 파일이 사용 중입니다: {output_path}")
            return False, "파일이 사용 중입니다. PDF 뷰어를 닫고 다시 시도하세요."
        
        # 2. 임시 파일명 생성
        temp_output = generate_safe_filename(f"temp_{os.path.basename(output_path)}")
        temp_output_path = os.path.join(os.path.dirname(output_path), temp_output)
        
        print(f"임시 파일 경로: {temp_output_path}")
        
        # 3. 변환 작업 수행
        try:
            success = docx_to_pdf(docx_path, temp_output_path)
            if not success:
                return False, "DOCX 변환 중 오류가 발생했습니다."
        except Exception as e:
            print(f"변환 오류: {e}")
            return False, f"변환 중 오류: {str(e)}"
        
        # 4. 임시 파일 확인
        if not os.path.exists(temp_output_path):
            return False, "변환된 파일이 생성되지 않았습니다."
        
        # 5. 임시 파일을 최종 파일로 이동
        try:
            if os.path.exists(output_path):
                os.remove(output_path)
            
            os.rename(temp_output_path, output_path)
            print(f"✅ 파일 이동 완료: {output_path}")
            
            return True, "변환이 성공적으로 완료되었습니다."
            
        except Exception as e:
            print(f"파일 이동 오류: {e}")
            return False, f"파일 이동 중 오류: {str(e)}"
            
    except Exception as e:
        print(f"❌ DOCX → PDF 변환 중 치명적 오류: {e}")
        return False, f"변환 중 오류가 발생했습니다: {str(e)}"

def pdf_to_docx(pdf_path, output_path, quality='medium'):
    """PDF를 DOCX로 변환"""
    try:
        print(f"🔄 PDF → DOCX 변환 시작: {pdf_path}")
        
        # 출력 파일이 이미 존재하면 삭제
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
                print(f"🗑️ 기존 파일 삭제: {output_path}")
            except Exception as e:
                print(f"⚠️ 기존 파일 삭제 실패: {e}")
        
        # PDF를 이미지로 변환
        quality_settings = {
            'medium': {'dpi': 150, 'format': 'jpeg'},
            'high': {'dpi': 300, 'format': 'png'}
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
        
        # 새 Word 문서 생성
        doc = Document()
        
        for i, image in enumerate(images):
            print(f"🔄 페이지 {i + 1}/{len(images)} 처리 중...")
            
            # 이미지를 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
                image.save(temp_file.name, 'JPEG', quality=85)
                
                # 문서에 이미지 추가
                doc.add_picture(temp_file.name, width=Inches(6))
                
                # 페이지 구분을 위한 페이지 브레이크 추가 (마지막 페이지 제외)
                if i < len(images) - 1:
                    doc.add_page_break()
            
            # 임시 파일 삭제
            try:
                os.unlink(temp_file.name)
            except:
                pass
        
        # DOCX 파일 저장
        doc.save(output_path)
        
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"✅ DOCX 변환 성공: {output_path} (크기: {file_size:,} bytes)")
            return True
        else:
            print(f"❌ DOCX 파일 생성 실패")
            return False
            
    except Exception as e:
        print(f"❌ PDF to DOCX 변환 오류: {e}")
        return False

def docx_to_pdf(docx_path, output_path):
    """DOCX를 PDF로 변환 (docx2pdf 없이)"""
    try:
        print(f"🔄 DOCX → PDF 변환 시작: {docx_path}")
        
        # python-docx + reportlab 사용
        from docx import Document
        
        doc = Document(docx_path)
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y_position = height - 50
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # 한글 폰트 설정 (나눔고딕)
                try:
                    font_path = os.path.join('fonts', 'NanumGothic.woff2')
                    if os.path.exists(font_path):
                        pdfmetrics.registerFont(TTFont('NanumGothic', font_path))
                        c.setFont('NanumGothic', 12)
                    else:
                        c.setFont('Helvetica', 12)
                except:
                    c.setFont('Helvetica', 12)
                
                # 텍스트 길이 제한 및 줄바꿈 처리
                text = paragraph.text[:100]
                c.drawString(50, y_position, text)
                y_position -= 20
                
                if y_position < 50:  # 페이지 끝에 도달하면 새 페이지
                    c.showPage()
                    y_position = height - 50
        
        c.save()
        
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"✅ DOCX → PDF 변환 성공: {output_path} (크기: {file_size:,} bytes)")
            return True
        else:
            print("❌ PDF 파일이 생성되지 않았습니다.")
            return False
            
    except Exception as e:
        print(f"❌ DOCX → PDF 변환 오류: {e}")
        return False

def create_pptx_table(slide, pdf_table, page_num, table_idx):
    """PDF 테이블을 PPTX 테이블로 변환"""
    try:
        # 테이블 데이터 추출
        table_data = pdf_table.extract()
        if not table_data:
            print(f"테이블 {table_idx+1}: 데이터 추출 실패")
            return
        
        # 빈 행 제거
        filtered_data = []
        for row in table_data:
            if row and any(cell and str(cell).strip() for cell in row):
                # None 값을 빈 문자열로 변환
                cleaned_row = [str(cell).strip() if cell else "" for cell in row]
                filtered_data.append(cleaned_row)
        
        if not filtered_data:
            print(f"테이블 {table_idx+1}: 유효한 데이터 없음")
            return
        
        rows = len(filtered_data)
        cols = max(len(row) for row in filtered_data) if filtered_data else 1
        
        print(f"테이블 {table_idx+1}: {rows}행 {cols}열 생성")
        
        # 테이블 위치 및 크기 계산
        bbox = pdf_table.bbox
        left = Inches(bbox[0] / 72 * 0.8)  # 약간 축소
        top = Inches(bbox[1] / 72 * 0.8)
        width = Inches((bbox[2] - bbox[0]) / 72 * 0.8)
        height = Inches((bbox[3] - bbox[1]) / 72 * 0.8)
        
        # 슬라이드 경계 확인
        if left.inches < 0: left = Inches(0.5)
        if top.inches < 0: top = Inches(0.5)
        if width.inches > 9: width = Inches(9)
        if height.inches > 6: height = Inches(6)
        
        # PPTX 테이블 생성
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # 테이블 데이터 입력 및 스타일 설정
        for row_idx, row_data in enumerate(filtered_data):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                if col_idx < len(row_data):
                    cell.text = str(row_data[col_idx])[:100]  # 텍스트 길이 제한
                else:
                    cell.text = ""
                
                # 셀 스타일 설정 (원본과 동일한 흑백 스타일)
                cell.text_frame.word_wrap = True
                cell.text_frame.margin_left = Inches(0.05)
                cell.text_frame.margin_right = Inches(0.05)
                cell.text_frame.margin_top = Inches(0.02)
                cell.text_frame.margin_bottom = Inches(0.02)
                
                # 셀 배경색을 흰색으로 설정
                from pptx.dml.color import RGBColor
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색 배경
                
                if cell.text_frame.paragraphs:
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(9)  # 원본에 맞는 작은 폰트
                    p.font.color.rgb = RGBColor(0, 0, 0)  # 검은색 텍스트
                    
                    # 헤더 행 감지 및 굵게 처리
                if row_idx == 0 or (cell.text and any(keyword in cell.text for keyword in ['신고인', '제목', '처리', '담당'])):
                    p.font.bold = True
        
        # 간단한 테이블 스타일 설정
        try:
            # 테이블에 기본 테두리 스타일 적용
            table_shape.table.style = None  # 기본 스타일 제거
            
            # 각 셀에 간단한 테두리 적용
            for row in table.rows:
                for cell in row.cells:
                    # 셀 배경을 흰색으로 유지
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    
        except Exception as style_error:
            print(f"테이블 스타일 설정 중 오류: {style_error}")
        
        print(f"테이블 {table_idx+1} 생성 완료")
        
    except Exception as e:
        print(f"테이블 생성 중 오류: {e}")
        # 오류 발생 시 텍스트로 대체
        try:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"테이블 {table_idx+1} 변환 실패 - 원본을 확인하세요"
            p.font.size = Pt(12)
        except:
            pass

def create_advanced_text_layout(slide, page, page_num):
    """고급 텍스트 레이아웃 분석 및 생성"""
    try:
        # 텍스트 객체들을 추출 (위치 정보 포함)
        chars = page.chars
        if not chars:
            # 기본 메시지
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"페이지 {page_num + 1} - 추출 가능한 텍스트가 없습니다"
            return
        
        # 텍스트를 블록별로 그룹화 (더 정교한 방식)
        text_blocks = group_chars_into_blocks(chars)
        
        for block_info in text_blocks:
            try:
                # 텍스트 박스 추가
                left = Inches(max(0, block_info['x0'] / 72))
                top = Inches(max(0, block_info['top'] / 72))
                width = Inches(min(9, (block_info['x1'] - block_info['x0']) / 72))
                height = Inches(min(6, block_info['height'] / 72))
                
                # 최소 크기 보장
                if width.inches < 0.5: width = Inches(0.5)
                if height.inches < 0.3: height = Inches(0.3)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.05)
                text_frame.margin_right = Inches(0.05)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)
                
                # 텍스트 추가
                p = text_frame.paragraphs[0]
                p.text = block_info['text'][:500]  # 텍스트 길이 제한
                
                # 폰트 크기 설정
                if block_info['size'] > 0:
                    p.font.size = Pt(min(max(block_info['size'], 8), 18))  # 8-18pt 범위
                else:
                    p.font.size = Pt(11)
                
                # 텍스트 정렬 (중앙 정렬된 텍스트 감지)
                if block_info.get('centered', False):
                    p.alignment = 1  # 중앙 정렬
                
            except Exception as e:
                print(f"텍스트 블록 처리 중 오류: {e}")
                continue
                
    except Exception as e:
        print(f"고급 텍스트 레이아웃 생성 중 오류: {e}")

def group_chars_into_blocks(chars):
    """문자들을 의미있는 텍스트 블록으로 그룹화"""
    if not chars:
        return []
    
    # Y 좌표와 X 좌표로 정렬
    sorted_chars = sorted(chars, key=lambda x: (x['top'], x['x0']))
    
    blocks = []
    current_block = []
    current_top = None
    tolerance_y = 5  # Y 좌표 허용 오차
    tolerance_x = 20  # X 좌표 허용 오차 (같은 라인 내)
    
    for char in sorted_chars:
        if current_top is None:
            current_block = [char]
            current_top = char['top']
        elif abs(char['top'] - current_top) <= tolerance_y:
            # 같은 라인으로 간주
            if current_block and abs(char['x0'] - current_block[-1]['x1']) <= tolerance_x:
                current_block.append(char)
            else:
                # 새로운 블록 시작
                if current_block:
                    blocks.append(create_block_info(current_block))
                current_block = [char]
        else:
            # 새로운 라인
            if current_block:
                blocks.append(create_block_info(current_block))
            current_block = [char]
            current_top = char['top']
    
    if current_block:
        blocks.append(create_block_info(current_block))
    
    return [block for block in blocks if block and block['text'].strip()]

def create_block_info(chars):
    """문자 리스트에서 블록 정보를 생성"""
    if not chars:
        return None
    
    # X 좌표로 정렬
    sorted_chars = sorted(chars, key=lambda x: x['x0'])
    
    text = ''.join([char['text'] for char in sorted_chars])
    x0 = min([char['x0'] for char in sorted_chars])
    x1 = max([char['x1'] for char in sorted_chars])
    top = min([char['top'] for char in sorted_chars])
    bottom = max([char['bottom'] for char in sorted_chars])
    
    # 평균 폰트 크기 계산
    sizes = [char.get('size', 12) for char in sorted_chars if char.get('size', 0) > 0]
    avg_size = sum(sizes) / len(sizes) if sizes else 12
    
    # 중앙 정렬 여부 감지 (간단한 휴리스틱)
    page_center = 300  # 대략적인 페이지 중앙
    text_center = (x0 + x1) / 2
    centered = abs(text_center - page_center) < 50
    
    return {
        'text': text.strip(),
        'x0': x0,
        'x1': x1,
        'top': top,
        'height': bottom - top,
        'size': avg_size,
        'centered': centered
    }

def group_chars_into_lines(chars):
    """문자들을 라인별로 그룹화하는 함수"""
    if not chars:
        return []
    
    # Y 좌표로 정렬
    sorted_chars = sorted(chars, key=lambda x: (x['top'], x['x0']))
    
    lines = []
    current_line = []
    current_top = None
    tolerance = 2  # Y 좌표 허용 오차
    
    for char in sorted_chars:
        if current_top is None or abs(char['top'] - current_top) <= tolerance:
            current_line.append(char)
            current_top = char['top']
        else:
            if current_line:
                lines.append(create_line_info(current_line))
            current_line = [char]
            current_top = char['top']
    
    if current_line:
        lines.append(create_line_info(current_line))
    
    return lines

def create_line_info(chars):
    """문자 리스트에서 라인 정보를 생성하는 함수"""
    if not chars:
        return None
    
    # X 좌표로 정렬
    sorted_chars = sorted(chars, key=lambda x: x['x0'])
    
    text = ''.join([char['text'] for char in sorted_chars])
    x0 = min([char['x0'] for char in sorted_chars])
    x1 = max([char['x1'] for char in sorted_chars])
    top = sorted_chars[0]['top']
    
    # 평균 폰트 크기 계산
    sizes = [char.get('size', 12) for char in sorted_chars if char.get('size', 0) > 0]
    avg_size = sum(sizes) / len(sizes) if sizes else 12
    
    return {
        'text': text.strip(),
        'x0': x0,
        'x1': x1,
        'top': top,
        'size': avg_size
    }

def pdf_to_pptx_with_text(pdf_path, output_path, quality='medium'):
    """PDF를 PPTX로 변환 (텍스트 계층 추출 우선)"""
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        print(f"📝 텍스트 계층 기반 PDF → PPTX 변환 시작: {pdf_path}")
        print(f"🎯 변환 방식: 내장 텍스트 추출 (OCR 없음)")
        
        # PDF 문서 열기
        pdf_document = fitz.open(pdf_path)
        prs = Presentation()
        
        # 전체 문서에서 추출된 텍스트를 저장 (중복 방지용)
        all_extracted_text = set()
        processed_pages = set()
        
        print(f"📄 총 {len(pdf_document)} 페이지 처리 시작")
        
        for page_num in range(len(pdf_document)):
            # 페이지 중복 처리 방지
            if page_num in processed_pages:
                print(f"⚠️ 페이지 {page_num + 1} 이미 처리됨 - 건너뛰기")
                continue
            
            try:
                page = pdf_document.load_page(page_num)
                
                # 페이지 텍스트 추출
                page_text = page.get_text()
                
                # 빈 페이지 건너뛰기
                if not page_text.strip():
                    print(f"📄 페이지 {page_num + 1}: 빈 페이지 건너뛰기")
                    processed_pages.add(page_num)
                    continue
                
                # 텍스트 정규화 (중복 체크용)
                normalized_text = ' '.join(page_text.split()).lower()
                
                # 중복 페이지 체크
                if normalized_text in all_extracted_text:
                    print(f"⚠️ 페이지 {page_num + 1}: 중복 내용 감지 - 건너뛰기")
                    processed_pages.add(page_num)
                    continue
                
                # 텍스트 정리
                cleaned_text = clean_extracted_text(page_text)
                
                if not cleaned_text.strip():
                    print(f"📄 페이지 {page_num + 1}: 정리 후 빈 페이지 - 건너뛰기")
                    processed_pages.add(page_num)
                    continue
                
                # 중복 방지를 위해 정규화된 텍스트 저장
                all_extracted_text.add(normalized_text)
                processed_pages.add(page_num)
                
                # 새 슬라이드 추가
                slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
                slide = prs.slides.add_slide(slide_layout)
                
                # ✅ 텍스트만 추가 (이미지 배경 없음)
                textbox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), 
                    Inches(9), Inches(6.5)
                )
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.2)
                text_frame.margin_right = Inches(0.2)
                text_frame.margin_top = Inches(0.2)
                text_frame.margin_bottom = Inches(0.2)
                
                # 텍스트 내용 추가
                lines = cleaned_text.split('\n')
                meaningful_lines = []
                
                for line in lines:
                    line = line.strip()
                    if line and len(line) > 2:
                        meaningful_lines.append(line)
                
                if meaningful_lines:
                    # 첫 번째 줄 (제목)
                    text_frame.text = meaningful_lines[0]
                    p = text_frame.paragraphs[0]
                    p.font.size = Pt(16)
                    p.font.name = 'NanumGothic'
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # 나머지 줄들 (본문)
                    for line in meaningful_lines[1:]:
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.font.size = Pt(12)
                        p.font.name = 'NanumGothic'
                        p.font.color.rgb = RGBColor(0, 0, 0)
                        p.space_after = Pt(6)
                
                print(f"✅ 페이지 {page_num + 1}: 텍스트 변환 완료 ({len(meaningful_lines)}줄)")
                
            except Exception as page_error:
                print(f"❌ 페이지 {page_num + 1} 처리 실패: {page_error}")
                processed_pages.add(page_num)
                continue
        
        # 빈 프레젠테이션 체크
        if len(prs.slides) == 0:
            print("⚠️ 변환 가능한 텍스트가 없어 빈 슬라이드 추가")
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            text_shape = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(4)
            )
            text_frame = text_shape.text_frame
            text_frame.text = "변환 가능한 텍스트 내용이 없습니다."
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # PPTX 파일 저장
        prs.save(output_path)
        pdf_document.close()
        
        print(f"🎉 텍스트 기반 변환 완료: {len(prs.slides)}개 슬라이드 생성")
        return True
        
    except Exception as e:
        print(f"❌ 텍스트 기반 변환 오류: {e}")
        return False

def check_pdf_text_extractable(pdf_path):
    """PDF에서 텍스트 추출이 가능한지 사전 검사"""
    try:
        import fitz  # PyMuPDF
        
        print(f"🔍 PDF 텍스트 계층 검사: {pdf_path}")
        
        # PDF 문서 열기
        pdf_document = fitz.open(pdf_path)
        
        # 처음 3페이지만 검사 (성능 최적화)
        pages_to_check = min(3, len(pdf_document))
        total_text_length = 0
        meaningful_text_count = 0
        
        for page_num in range(pages_to_check):
            try:
                page = pdf_document.load_page(page_num)
                page_text = page.get_text().strip()
                
                if page_text:
                    total_text_length += len(page_text)
                    
                    # 의미있는 텍스트인지 확인
                    words = page_text.split()
                    meaningful_words = [w for w in words if len(w) > 1 and not w.isdigit()]
                    
                    if len(meaningful_words) > 5:  # 의미있는 단어가 5개 이상
                        meaningful_text_count += 1
                        
            except Exception as e:
                print(f"⚠️ 페이지 {page_num + 1} 텍스트 검사 오류: {e}")
                continue
        
        pdf_document.close()
        
        # 판정 기준
        avg_text_per_page = total_text_length / pages_to_check if pages_to_check > 0 else 0
        meaningful_ratio = meaningful_text_count / pages_to_check if pages_to_check > 0 else 0
        
        print(f"📊 텍스트 분석 결과:")
        print(f"   - 검사 페이지: {pages_to_check}페이지")
        print(f"   - 평균 텍스트 길이: {avg_text_per_page:.0f}자")
        print(f"   - 의미있는 페이지 비율: {meaningful_ratio:.1%}")
        
        # 텍스트 추출 가능 판정
        if avg_text_per_page > 100 and meaningful_ratio > 0.5:
            print("✅ 텍스트 계층 추출 가능 (일반 PDF)")
            return True
        else:
            print("❌ 텍스트 계층 부족 (스캔 PDF 가능성)")
            return False
            
    except ImportError:
        print("⚠️ PyMuPDF 모듈 없음 - OCR로 진행")
        return False
    except Exception as e:
        print(f"⚠️ PDF 텍스트 검사 오류: {e} - OCR로 진행")
        return False

def extract_document_numbers(text):
    """문서에서 번호 필드 정교 추출"""
    extracted_numbers = {}
    
    # 🔥 한국 문서 번호 패턴들
    patterns = {
        'registration_number': [
            r'등록번호[:\s]*([A-Z0-9-]+)',
            r'신고번호[:\s]*([A-Z0-9-]+)',
            r'승인번호[:\s]*([A-Z0-9-]+)',
            r'관리번호[:\s]*([A-Z0-9-]+)'
        ],
        'business_number': [
            r'사업자등록번호[:\s]*([0-9-]+)',
            r'사업자번호[:\s]*([0-9-]+)'
        ],
        'document_number': [
            r'문서번호[:\s]*([A-Z0-9-]+)',
            r'문서[\s]*번호[:\s]*([A-Z0-9-]+)'
        ],
        'kc_number': [
            r'KC[\s]*인증[\s]*번호[:\s]*([A-Z0-9-]+)',
            r'KC[:\s]*([A-Z0-9-]+)',
            r'안전확인[\s]*번호[:\s]*([A-Z0-9-]+)'
        ],
        'phone_number': [
            r'전화번호?[:\s]*([0-9-]+)',
            r'연락처[:\s]*([0-9-]+)',
            r'TEL[:\s]*([0-9-]+)'
        ],
        'date': [
            r'발행일자?[:\s]*([0-9]{4}[.-][0-9]{1,2}[.-][0-9]{1,2})',
            r'유효기간[:\s]*([0-9]{4}[.-][0-9]{1,2}[.-][0-9]{1,2})'
        ]
    }
    
    for field_name, pattern_list in patterns.items():
        for pattern in pattern_list:
            matches = re.findall(pattern, text, re.IGNORECASE)
            if matches:
                extracted_numbers[field_name] = matches[0].strip()
                print(f"📋 {field_name}: {matches[0].strip()}")
                break
    
    return extracted_numbers

def pdf_to_pptx(pdf_path, output_path, quality='medium'):
    """PDF → PPTX 변환 (텍스트 우선 + 번호 추출)"""
    start_time = time.time()
    
    try:
        print(f"🔄 PDF → PPTX 변환 시작: {pdf_path}")
        
        # 출력 파일 정리
        if os.path.exists(output_path):
            os.remove(output_path)
        
        success = False
        conversion_method = ""
        extracted_numbers = {}
        
        # 🔥 1단계: 텍스트 계층 추출 시도
        print("📝 1단계: PDF 텍스트 계층 추출 시도")
        text_extractable = check_pdf_text_extractable(pdf_path)
        
        if text_extractable:
            print("✅ 텍스트 계층 존재 - 직접 추출 진행")
            
            # 텍스트 추출 및 번호 필드 분석
            try:
                import fitz
                pdf_doc = fitz.open(pdf_path)
                full_text = ""
                
                for page_num in range(min(5, len(pdf_doc))):
                    page = pdf_doc.load_page(page_num)
                    full_text += page.get_text() + "\n"
                
                pdf_doc.close()
                
                # 🔥 번호 필드 추출
                extracted_numbers = extract_document_numbers(full_text)
                print(f"📋 추출된 번호 필드: {len(extracted_numbers)}개")
                
            except Exception as e:
                print(f"⚠️ 번호 추출 오류: {e}")
            
            # 텍스트 기반 변환
            success = pdf_to_pptx_with_text(pdf_path, output_path, quality)
            conversion_method = "텍스트 계층 추출"
            
            if success:
                print(f"🎉 텍스트 추출 성공 - OCR 건너뛰기")
                
                # 🔥 파일명 재구성
                if extracted_numbers:
                    base_dir = os.path.dirname(output_path)
                    original_filename = os.path.basename(output_path)
                    new_filename = generate_safe_filename(original_filename, extracted_numbers)
                    new_output_path = os.path.join(base_dir, new_filename)
                    
                    try:
                        os.rename(output_path, new_output_path)
                        print(f"📝 파일명 변경: {new_filename}")
                        output_path = new_output_path
                    except Exception as e:
                        print(f"⚠️ 파일명 변경 실패: {e}")
                
                # 변환 완료 후 DB 저장
                processing_time = time.time() - start_time
                doc_id = doc_manager.save_document_data(
                    pdf_path=pdf_path,
                    extracted_numbers=extracted_numbers,
                    conversion_method=conversion_method,
                    success=True,
                    processing_time=processing_time
                )
                print(f"💾 문서 데이터 저장 완료: ID {doc_id}")
                
                return True
        
        # 🔥 2단계: OCR 폴백
        if not success and OCR_AVAILABLE:
            print("🔍 2단계: OCR 기반 변환 시도")
            success = pdf_to_pptx_with_ocr(pdf_path, output_path, quality)
            conversion_method = "OCR 기반"
            
            if success:
                print(f"✅ OCR 변환 성공")
                
                # 변환 완료 후 DB 저장
                processing_time = time.time() - start_time
                doc_id = doc_manager.save_document_data(
                    pdf_path=pdf_path,
                    extracted_numbers=extracted_numbers,
                    conversion_method=conversion_method,
                    success=True,
                    processing_time=processing_time
                )
                print(f"💾 문서 데이터 저장 완료: ID {doc_id}")
                
                return True
        
        # 🔥 3단계: 이미지 기반 폴백
        if not success:
            print("🖼️ 3단계: 이미지 기반 변환 (최종 폴백)")
            success = pdf_to_pptx_with_images_only(pdf_path, output_path, quality)
            conversion_method = "이미지 기반"
        
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"✅ {conversion_method} 변환 성공: {file_size:,} bytes")
            
            # 변환 완료 후 DB 저장
            processing_time = time.time() - start_time
            doc_id = doc_manager.save_document_data(
                pdf_path=pdf_path,
                extracted_numbers=extracted_numbers,
                conversion_method=conversion_method,
                success=True,
                processing_time=processing_time
            )
            print(f"💾 문서 데이터 저장 완료: ID {doc_id}")
            
            return True
        else:
            print(f"❌ 모든 변환 방식 실패")
            
            # 실패 정보 DB 저장
            processing_time = time.time() - start_time
            doc_manager.save_document_data(
                pdf_path=pdf_path,
                extracted_numbers={},
                conversion_method="failed",
                success=False,
                processing_time=processing_time
            )
            
            return False
            
    except Exception as e:
        processing_time = time.time() - start_time
        doc_manager.save_document_data(
            pdf_path=pdf_path,
            extracted_numbers={},
            conversion_method="error",
            success=False,
            processing_time=processing_time
        )
        print(f"❌ 변환 오류: {e}")
        return False

def pptx_to_pdf(pptx_path, output_path):
    """PPTX를 PDF로 변환 (슬라이드 제목 제거 + 나눔고딕 폰트 강화)"""
    try:
        print(f"🔄 PPTX → PDF 변환 시작: {pptx_path}")
        
        # 1단계: 입력 파일 검증
        if not os.path.exists(pptx_path):
            print(f"❌ 입력 파일이 존재하지 않습니다: {pptx_path}")
            return False
        
        if not pptx_path.lower().endswith(('.pptx', '.ppt')):
            print(f"❌ 지원하지 않는 파일 형식: {pptx_path}")
            return False
        
        # 2단계: 출력 디렉토리 생성
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir, exist_ok=True)
                print(f"📁 출력 디렉토리 생성: {output_dir}")
            except Exception as e:
                print(f"❌ 출력 디렉토리 생성 실패: {e}")
                return False
        
        # 3단계: 나눔고딕 폰트 자동 다운로드 및 등록 (강화된 버전)
        korean_font = None
        
        def download_korean_font():
            """나눔고딕 폰트 자동 다운로드"""
            try:
                import urllib.request
                import zipfile
                
                # fonts 디렉토리 생성
                fonts_dir = os.path.join(os.getcwd(), 'fonts')
                os.makedirs(fonts_dir, exist_ok=True)
                
                # 나눔고딕 폰트 다운로드 URL
                font_url = "https://github.com/naver/nanumfont/releases/download/VER2.6/NanumFont_TTF_ALL.zip"
                font_zip_path = os.path.join(fonts_dir, "NanumFont.zip")
                font_ttf_path = os.path.join(fonts_dir, "NanumGothic.ttf")
                
                # 이미 폰트가 있으면 건너뛰기
                if os.path.exists(font_ttf_path):
                    print(f"✅ 기존 나눔고딕 폰트 발견: {font_ttf_path}")
                    return font_ttf_path
                
                print(f"📥 나눔고딕 폰트 다운로드 중: {font_url}")
                
                # 폰트 다운로드
                urllib.request.urlretrieve(font_url, font_zip_path)
                print(f"✅ 폰트 다운로드 완료: {font_zip_path}")
                
                # ZIP 파일 압축 해제
                with zipfile.ZipFile(font_zip_path, 'r') as zip_ref:
                    # NanumGothic.ttf 파일만 추출
                    for file_info in zip_ref.filelist:
                        if file_info.filename.endswith('NanumGothic.ttf'):
                            file_info.filename = 'NanumGothic.ttf'  # 파일명 단순화
                            zip_ref.extract(file_info, fonts_dir)
                            break
                
                # ZIP 파일 삭제
                os.remove(font_zip_path)
                
                if os.path.exists(font_ttf_path):
                    print(f"✅ 나눔고딕 폰트 설치 완료: {font_ttf_path}")
                    return font_ttf_path
                else:
                    print(f"❌ 폰트 추출 실패")
                    return None
                    
            except Exception as e:
                print(f"❌ 폰트 다운로드 실패: {e}")
                return None
        
        # 나눔고딕 폰트 경로 목록 (우선순위: 나눔고딕 최우선)
        font_paths = [
            # 로컬 다운로드 나눔고딕 (최우선)
            os.path.join(os.getcwd(), 'fonts', 'NanumGothic.ttf'),
            
            # Windows 나눔고딕
            "C:/Windows/Fonts/NanumGothic.ttf",
            "C:/Windows/Fonts/NanumGothicBold.ttf",
            
            # Windows 기타 한글 폰트
            "C:/Windows/Fonts/malgun.ttf",      # 맑은 고딕
            "C:/Windows/Fonts/gulim.ttc",       # 굴림
            "C:/Windows/Fonts/batang.ttc",      # 바탕
            "C:/Windows/Fonts/dotum.ttc",       # 돋움
            
            # macOS 나눔고딕
            "/Library/Fonts/NanumGothic.ttf",
            "/System/Library/Fonts/AppleGothic.ttf",
            
            # Linux 나눔고딕
            "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
        ]
        
        # 폰트 등록 시도 (나눔고딕 우선)
        for font_path in font_paths:
            if os.path.exists(font_path):
                try:
                    pdfmetrics.registerFont(TTFont('NanumGothic', font_path))
                    korean_font = 'NanumGothic'
                    print(f"✅ 나눔고딕 폰트 등록 성공: {font_path}")
                    break
                except Exception as font_error:
                    print(f"⚠️ 폰트 등록 실패: {font_path} - {font_error}")
                    continue
        
        # 폰트가 없으면 자동 다운로드
        if korean_font is None:
            print(f"⚠️ 시스템에서 나눔고딕 폰트를 찾을 수 없습니다. 자동 다운로드를 시도합니다...")
            downloaded_font = download_korean_font()
            
            if downloaded_font and os.path.exists(downloaded_font):
                try:
                    pdfmetrics.registerFont(TTFont('NanumGothic', downloaded_font))
                    korean_font = 'NanumGothic'
                    print(f"✅ 다운로드된 나눔고딕 폰트 등록 성공: {downloaded_font}")
                except Exception as e:
                    print(f"❌ 다운로드된 폰트 등록 실패: {e}")
        
        # 최종적으로 폰트가 없으면 기본 폰트 사용
        if korean_font is None:
            korean_font = 'Helvetica'
            print(f"⚠️ 나눔고딕 폰트를 사용할 수 없어 기본 폰트를 사용합니다.")
        
        # 4단계: PPTX 파일 열기
        try:
            prs = Presentation(pptx_path)
            print(f"📄 PPTX 파일 로드 완료: {len(prs.slides)}개 슬라이드")
        except Exception as e:
            print(f"❌ PPTX 파일 열기 실패: {e}")
            return False
        
        # 5단계: PDF 캔버스 생성 (나눔고딕 지원 설정)
        try:
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            print(f"📄 PDF 캔버스 생성 완료: {width}x{height}")
        except Exception as e:
            print(f"❌ PDF 캔버스 생성 실패: {e}")
            return False
        
        # 6단계: 슬라이드 변환 (제목 제거 + 나눔고딕 적용)
        converted_slides = 0
        
        for slide_num, slide in enumerate(prs.slides):
            try:
                print(f"🔄 슬라이드 {slide_num + 1}/{len(prs.slides)} 처리 중...")
                
                # 새 페이지 시작 (첫 번째 슬라이드 제외)
                if slide_num > 0:
                    c.showPage()
                
                # ❌ 페이지 제목 추가 부분 완전 제거 (사용자 요청)
                # 제목 없이 바로 내용부터 시작
                
                y_position = height - 50  # 제목 공간 제거하여 위쪽부터 시작
                
                # 7단계: 텍스트 처리 (나눔고딕 폰트 강화)
                text_shapes_processed = 0
                
                for shape_num, shape in enumerate(slide.shapes):
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            # 텍스트 내용 가져오기 및 인코딩 처리
                            text_content = shape.text.strip()
                            
                            # UTF-8 인코딩 확인 및 처리
                            try:
                                text_content.encode('utf-8')
                            except UnicodeEncodeError:
                                text_content = text_content.encode('utf-8', errors='replace').decode('utf-8')
                            
                            if not text_content:
                                continue
                            
                            # 나눔고딕 폰트 크기 설정
                            try:
                                if hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 0:
                                    c.setFont(korean_font, 16)  # 제목용 크기
                                    font_size = 16
                                else:
                                    c.setFont(korean_font, 12)  # 본문용 크기
                                    font_size = 12
                                print(f"✅ 나눔고딕 폰트 적용: {korean_font}, 크기: {font_size}")
                            except Exception as font_error:
                                print(f"⚠️ 나눔고딕 폰트 설정 실패: {font_error}")
                                try:
                                    c.setFont('Helvetica', 12)
                                    font_size = 12
                                except:
                                    font_size = 12
                            
                            # 텍스트 줄바꿈 처리 (나눔고딕 한글 지원)
                            text_lines = text_content.split('\n')
                            
                            for line in text_lines:
                                line = line.strip()
                                if line and y_position > 50:
                                    try:
                                        # 나눔고딕 한글 문자 폭 계산
                                        korean_char_count = sum(1 for char in line if ord(char) > 127)
                                        english_char_count = len(line) - korean_char_count
                                        
                                        # 나눔고딕 폰트의 한글과 영문 폭 계산
                                        estimated_width = (korean_char_count * font_size * 0.9) + (english_char_count * font_size * 0.6)
                                        max_width = width - 100
                                        
                                        if estimated_width > max_width:
                                            # 긴 텍스트 줄바꿈 처리
                                            words = line.split(' ')
                                            current_line = ""
                                            
                                            for word in words:
                                                test_line = current_line + " " + word if current_line else word
                                                
                                                # 테스트 라인의 예상 폭 계산
                                                test_korean = sum(1 for char in test_line if ord(char) > 127)
                                                test_english = len(test_line) - test_korean
                                                test_width = (test_korean * font_size * 0.9) + (test_english * font_size * 0.6)
                                                
                                                if test_width <= max_width:
                                                    current_line = test_line
                                                else:
                                                    if current_line:
                                                        # 현재 줄 출력 (나눔고딕)
                                                        c.drawString(50, y_position, current_line)
                                                        y_position -= (font_size + 4)
                                                        if y_position <= 50:
                                                            break
                                                    current_line = word
                                            
                                            # 마지막 줄 처리
                                            if current_line and y_position > 50:
                                                c.drawString(50, y_position, current_line)
                                                y_position -= (font_size + 4)
                                        else:
                                            # 짧은 텍스트는 그대로 출력 (나눔고딕)
                                            c.drawString(50, y_position, line)
                                            y_position -= (font_size + 4)
                                        
                                        print(f"✅ 나눔고딕 텍스트 출력: {line[:30]}{'...' if len(line) > 30 else ''}")
                                    
                                    except Exception as text_error:
                                        print(f"⚠️ 텍스트 출력 실패: {text_error}")
                                        # 안전한 폴백 처리
                                        try:
                                            c.setFont('Helvetica', 11)
                                            safe_text = line.encode('ascii', errors='ignore').decode('ascii')
                                            if safe_text.strip():
                                                c.drawString(50, y_position, safe_text[:50])
                                            y_position -= 15
                                        except:
                                            pass
                            
                            y_position -= 10  # 도형 간 간격
                            text_shapes_processed += 1
                    
                    except Exception as shape_error:
                        print(f"⚠️ 도형 {shape_num + 1} 처리 실패: {shape_error}")
                        continue
                
                # 페이지 번호 추가 (나눔고딕 지원)
                try:
                    c.setFont(korean_font, 9)
                    page_text = f"{slide_num + 1}"
                    c.drawString(width - 30, 20, page_text)
                except:
                    try:
                        c.setFont('Helvetica', 9)
                        c.drawString(width - 30, 20, f"{slide_num + 1}")
                    except:
                        pass
                
                converted_slides += 1
                print(f"✅ 슬라이드 {slide_num + 1} 변환 완료 (텍스트 도형: {text_shapes_processed}개)")
                
            except Exception as slide_error:
                print(f"❌ 슬라이드 {slide_num + 1} 처리 실패: {slide_error}")
                continue
        
        # 8단계: PDF 저장
        try:
            c.save()
            print(f"💾 PDF 저장 완료: {output_path}")
        except Exception as save_error:
            print(f"❌ PDF 저장 실패: {save_error}")
            return False
        
        # 9단계: 결과 검증
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"🎉 나눔고딕 변환 성공: {converted_slides}/{len(prs.slides)} 슬라이드 변환 완료")
            print(f"📊 파일 크기: {file_size:,} bytes")
            print(f"🔤 사용된 폰트: {korean_font}")
            return True
        else:
            print(f"❌ 출력 파일이 생성되지 않았습니다: {output_path}")
            return False
        
    except Exception as e:
        print(f"❌ PPTX → PDF 변환 중 치명적 오류: {e}")
        print(f"📍 오류 위치: {traceback.format_exc()}")
        return False
    
    finally:
        # 메모리 정리
        try:
            if 'prs' in locals():
                del prs
            if 'c' in locals():
                del c
        except:
            pass

# 파일 크기 초과 오류 처리
@app.errorhandler(413)
def too_large(e):
    return jsonify({'success': False, 'error': '파일 크기가 너무 큽니다. 100MB 이하의 파일을 업로드하세요.'}), 413

@app.route('/') 
def index(): 
    return render_template('index.html') 

@app.route('/ocr_review')
def ocr_review():
    """OCR 결과 검수 페이지"""
    return render_template('ocr_review.html')

@app.route('/save_review', methods=['POST'])
def save_review():
    """검수된 OCR 결과 저장"""
    try:
        data = request.get_json()
        reviewed_data = data.get('reviewedData', [])
        
        # 검수된 데이터를 파일로 저장하거나 DB에 저장
        import json
        with open('outputs/reviewed_ocr_data.json', 'w', encoding='utf-8') as f:
            json.dump(reviewed_data, f, ensure_ascii=False, indent=2)
        
        return jsonify({'success': True, 'message': '검수 데이터가 저장되었습니다.'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/convert', methods=['POST'])
def convert_file():
    """파일 변환 처리 (404 오류 수정)"""
    try:
        print("=== 변환 요청 시작 ===")
        
        # 1단계: 파일 존재 여부 확인
        if 'file' not in request.files:
            print("1단계 실패: 파일이 선택되지 않음")
            return jsonify({'success': False, 'error': '파일이 선택되지 않았습니다.'}), 400
        
        file = request.files['file']
        
        # 2단계: 파일명 확인
        if file.filename == '':
            print("2단계 실패: 파일명이 비어있음")
            return jsonify({'success': False, 'error': '파일명이 비어있습니다.'}), 400
        
        # 파일 크기 확인
        file.seek(0, 2)
        file_size = file.tell()
        file.seek(0)
        print(f"파일 크기: {file_size/1024/1024:.1f}MB")
        
        # 3단계: 파일 크기 확인 (100MB 제한)
        if file_size > 100 * 1024 * 1024:
            print(f"3단계 실패: 파일 크기 초과 ({file_size/1024/1024:.1f}MB)")
            return jsonify({
                'success': False, 
                'error': '파일 크기가 100MB를 초과합니다.',
                'file_size': f"{file_size/1024/1024:.1f}MB"
            }), 413
        
        # 4단계: 파일 형식 확인
        if not allowed_file(file.filename):
            print("4단계 실패: 지원하지 않는 파일 형식")
            return jsonify({
                'success': False, 
                'error': '지원하지 않는 파일 형식입니다. PDF 또는 DOCX 파일만 업로드 가능합니다.',
                'filename': file.filename
            }), 400
        
        # 5단계: 파일 저장
        filename = secure_filename(file.filename)
        if not filename:
            filename = f"upload_{int(time.time())}.{file.filename.rsplit('.', 1)[1].lower()}"
        
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        print(f"5단계: 파일 저장 중 - {file_path}")
        
        try:
            file.save(file_path)
            print("파일 저장 완료")
        except Exception as e:
            print(f"파일 저장 실패: {e}")
            return jsonify({
                'success': False, 
                'error': f'파일 저장 중 오류가 발생했습니다: {str(e)}'
            }), 500
        
        # 6단계: 변환 설정 확인
        quality = request.form.get('quality', 'medium')
        print(f"6단계: 변환 설정 - 품질: {quality}")
        
        # 7단계: 변환 시작
        file_ext = filename.rsplit('.', 1)[1].lower()
        print(f"파일 확장자: {file_ext}")
        
        if file_ext == 'pdf':
            # PDF → DOCX 변환
            output_filename = filename.rsplit('.', 1)[0] + '.docx'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            print(f"7단계: PDF→DOCX 변환 시작 - {file_path} -> {output_path}")
            
            success, message = pdf_to_docx_safe(file_path, output_path, quality)
            
        elif file_ext == 'docx':
            # DOCX → PDF 변환
            output_filename = filename.rsplit('.', 1)[0] + '.pdf'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            print(f"7단계: DOCX→PDF 변환 시작 - {file_path} -> {output_path}")
            
            success, message = docx_to_pdf_safe(file_path, output_path)
        else:
            return jsonify({
                'success': False, 
                'error': f'지원하지 않는 파일 확장자: {file_ext}'
            }), 400
        
        # 8단계: 변환 결과 처리
        if success:
            print("8단계: 변환 성공 - 다운로드 준비")
            
            # 임시 파일 삭제
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    print("임시 파일 삭제 완료")
            except Exception as e:
                print(f"임시 파일 삭제 실패: {e}")
            
            # 출력 파일 존재 확인
            if not os.path.exists(output_path):
                return jsonify({
                    'success': False, 
                    'error': '변환된 파일을 찾을 수 없습니다.',
                    'message': message
                }), 500
            
            # 9단계: 파일 다운로드
            print("9단계: 파일 다운로드 시작")
            try:
                response = send_file(output_path, as_attachment=True, download_name=output_filename)
                
                # CORS 헤더 추가
                response.headers['Access-Control-Allow-Origin'] = '*'
                response.headers['Access-Control-Allow-Methods'] = 'POST, GET, OPTIONS'
                response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
                
                return response
                
            except Exception as e:
                print(f"파일 다운로드 오류: {e}")
                return jsonify({
                    'success': False, 
                    'error': f'파일 다운로드 중 오류가 발생했습니다: {str(e)}'
                }), 500
        else:
            print(f"8단계: 변환 실패 - {message}")
            
            # 실패한 파일들 정리
            for cleanup_path in [file_path, output_path]:
                try:
                    if cleanup_path and os.path.exists(cleanup_path):
                        os.remove(cleanup_path)
                except Exception as e:
                    print(f"파일 정리 실패: {e}")
            
            return jsonify({
                'success': False, 
                'error': message or '파일 변환에 실패했습니다.'
            }), 500
            
    except Exception as e:
        print(f"❌ 변환 처리 중 치명적 오류: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False, 
            'error': f'서버 오류가 발생했습니다: {str(e)}'
        }), 500

if __name__ == '__main__':
    print("🚀 개선된 PDF ↔ DOCX 변환기 시작")
    print("✅ 안정적인 텍스트 기반 변환")
    print("✅ 한글 인코딩 지원")
    print("✅ DOCX 변환 완료")
    app.run(debug=True, host='0.0.0.0', port=5001)  # 포트 변경