from flask import Flask, request, render_template, send_file, flash, redirect, url_for
import os
import tempfile
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import json
from dotenv import load_dotenv
# Adobe PDF Services SDK imports (선택적 - API 키가 설정된 경우에만 사용)
try:
    from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
    from adobe.pdfservices.operation.exception.exceptions import ServiceApiException, ServiceUsageException, SdkException
    from adobe.pdfservices.operation.pdf_services import PDFServices
    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
    from adobe.pdfservices.operation.io.cloud_asset import CloudAsset
    from adobe.pdfservices.operation.io.stream_asset import StreamAsset
    ADOBE_SDK_AVAILABLE = True
except ImportError:
    print("Adobe PDF Services SDK를 사용할 수 없습니다. 기본 변환 방식을 사용합니다.")
    ADOBE_SDK_AVAILABLE = False

# 환경 변수 로드
load_dotenv()

app = Flask(__name__)
app.secret_key = 'your-secret-key-here'
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size

# Adobe PDF Services API 설정
ADOBE_CONFIG = {
    "client_credentials": {
        "client_id": os.getenv("ADOBE_CLIENT_ID", "YOUR_CLIENT_ID"),
        "client_secret": os.getenv("ADOBE_CLIENT_SECRET", "YOUR_CLIENT_SECRET")
    },
    "service_principal_credentials": {
        "organization_id": os.getenv("ADOBE_ORGANIZATION_ID", "3C67227E688C66000A495C72@AdobeOrg")
    }
}

UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ALLOWED_EXTENSIONS = {'pdf'}

# 폴더 생성
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_pdf_content_with_adobe(pdf_path):
    """Adobe PDF Services API를 사용하여 PDF 내용을 추출하는 함수"""
    if not ADOBE_SDK_AVAILABLE:
        print("Adobe PDF Services SDK를 사용할 수 없습니다.")
        return None
        
    try:
        # Adobe API 자격 증명 설정
        credentials = ServicePrincipalCredentials(
            client_id=ADOBE_CONFIG["client_credentials"]["client_id"],
            client_secret=ADOBE_CONFIG["client_credentials"]["client_secret"],
            organization_id=ADOBE_CONFIG["service_principal_credentials"]["organization_id"]
        )
        
        # PDF Services 인스턴스 생성
        pdf_services = PDFServices(credentials=credentials)
        
        # PDF 파일을 스트림으로 읽기
        with open(pdf_path, 'rb') as file:
            input_stream = file.read()
        
        # StreamAsset 생성
        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
        
        print("Adobe API를 사용하여 PDF 내용을 처리했습니다.")
        return input_asset
            
    except (ServiceApiException, ServiceUsageException, SdkException) as e:
        print(f"Adobe API 오류: {str(e)}")
        return None
    except Exception as e:
        print(f"일반 오류: {str(e)}")
        return None

def pdf_to_docx(pdf_path, output_path, quality='medium'):
    """PDF를 DOCX로 변환하는 함수 (Adobe API 통합)"""
    try:
        # 품질 설정에 따른 파라미터 설정 (최적화됨)
        quality_settings = {
            'medium': {
                'dpi': 120,  # DPI 최적화로 속도 향상
                'format': 'jpeg',
                'jpeg_quality': 80,  # 품질과 속도의 균형
                'max_size': (1600, 1200),  # 적절한 해상도
                'description': '균형 변환 (최적화된 속도와 품질)'
            },
            'high': {
                'dpi': 180,  # 고품질이지만 속도 고려
                'format': 'jpeg',  # PNG 대신 JPEG 사용으로 속도 향상
                'jpeg_quality': 90,
                'max_size': (2048, 1536),  # 해상도 최적화
                'description': '고품질 변환 (향상된 속도)'
            }
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        print(f"변환 설정: {settings['description']}")
        
        # Adobe API를 사용하여 PDF 내용 추출 시도
        if ADOBE_CONFIG["client_credentials"]["client_id"] != "YOUR_CLIENT_ID":
            extracted_content = extract_pdf_content_with_adobe(pdf_path)
            if extracted_content:
                print("Adobe API를 사용하여 PDF 내용을 추출했습니다.")
                # 추출된 내용을 기반으로 DOCX 생성 (향후 구현 가능)
        
        # 기본 방법: PDF를 이미지로 변환 (품질별 최적화)
        print("PDF를 이미지로 변환 중...")
        images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
        
        # 새 Word 문서 생성
        doc = Document()
        
        print(f"총 {len(images)}페이지 처리 중...")
        # 각 페이지를 문서에 추가
        for i, image in enumerate(images):
            print(f"페이지 {i+1}/{len(images)} 처리 중...")
            
            # 이미지 크기 최적화 (품질별 설정)
            max_width, max_height = settings['max_size']
            if image.size[0] > max_width or image.size[1] > max_height:
                image.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            
            # 이미지를 임시 파일로 저장 (JPEG 최적화)
            temp_img_path = None
            try:
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                    temp_img_path = temp_img.name
                    # JPEG로 저장 (품질별 압축, 빠른 처리)
                    image.save(temp_img_path, 'JPEG', quality=settings['jpeg_quality'], optimize=True)
                
                # 문서에 이미지 추가
                doc.add_picture(temp_img_path, width=Inches(6))
                
                # 페이지 구분을 위한 페이지 브레이크 추가 (마지막 페이지 제외)
                if i < len(images) - 1:
                    doc.add_page_break()
                
            finally:
                # 임시 파일 삭제 (빠른 처리)
                if temp_img_path and os.path.exists(temp_img_path):
                    try:
                        os.unlink(temp_img_path)
                    except (OSError, PermissionError) as e:
                        print(f"임시 파일 삭제 실패 (무시됨): {e}")
                        # 임시 파일 삭제 실패는 무시하고 계속 진행
        
        # DOCX 파일 저장
        doc.save(output_path)
        return True
        
    except Exception as e:
        print(f"변환 중 오류 발생: {str(e)}")
        return False

def pdf_to_pptx(pdf_path, output_path, quality='medium'):
    """PDF를 PPTX로 변환하는 함수 (Adobe API 통합)"""
    try:
        # 품질 설정에 따른 파라미터 설정 (최적화됨)
        quality_settings = {
            'medium': {
                'dpi': 120,  # DPI 최적화로 속도 향상
                'format': 'jpeg',
                'jpeg_quality': 80,  # 품질과 속도의 균형
                'max_size': (1600, 1200),  # 적절한 해상도
                'description': '균형 변환 (최적화된 속도와 품질)'
            },
            'high': {
                'dpi': 180,  # 고품질이지만 속도 고려
                'format': 'jpeg',  # PNG 대신 JPEG 사용으로 속도 향상
                'jpeg_quality': 90,
                'max_size': (2048, 1536),  # 해상도 최적화
                'description': '고품질 변환 (향상된 속도)'
            }
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        print(f"변환 설정: {settings['description']}")
        
        # Adobe API를 사용하여 PDF 내용 추출 시도
        if ADOBE_CONFIG["client_credentials"]["client_id"] != "YOUR_CLIENT_ID":
            extracted_content = extract_pdf_content_with_adobe(pdf_path)
            if extracted_content:
                print("Adobe API를 사용하여 PDF 내용을 추출했습니다.")
                # 추출된 내용을 기반으로 PPTX 생성 (향후 구현 가능)
        
        # 기본 방법: PDF를 이미지로 변환 (품질별 최적화)
        print("PDF를 이미지로 변환 중...")
        images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
        
        # 새 PowerPoint 프레젠테이션 생성
        prs = Presentation()
        
        print(f"총 {len(images)}페이지 처리 중...")
        def get_blank_slide_layout(prs):
            """안전한 빈 슬라이드 레이아웃 가져오기"""
            try:
                if len(prs.slide_layouts) > 6:
                    return prs.slide_layouts[6]  # 빈 슬라이드
                elif len(prs.slide_layouts) > 5:
                    return prs.slide_layouts[5]  # 제목만 있는 슬라이드
                elif len(prs.slide_layouts) > 0:
                    return prs.slide_layouts[0]  # 첫 번째 사용 가능한 레이아웃
                else:
                    return prs.slide_layouts[0]
            except IndexError:
                return prs.slide_layouts[0]
        
        # 각 페이지를 슬라이드로 추가
        for i, image in enumerate(images):
            print(f"페이지 {i+1}/{len(images)} 처리 중...")
            
            # 슬라이드 추가 - 안전한 레이아웃 사용
            slide_layout = get_blank_slide_layout(prs)
            slide = prs.slides.add_slide(slide_layout)
            
            # 이미지 크기 최적화 (품질별 설정)
            max_width, max_height = settings['max_size']
            if image.size[0] > max_width or image.size[1] > max_height:
                image.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            
            # 이미지를 임시 파일로 저장 (JPEG 최적화)
            temp_img_path = None
            try:
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                    temp_img_path = temp_img.name
                    # JPEG로 저장 (품질별 압축, 빠른 처리)
                    image.save(temp_img_path, 'JPEG', quality=settings['jpeg_quality'], optimize=True)
                
                # 슬라이드에 이미지 추가
                left = Inches(0.5)
                top = Inches(0.5)
                height = Inches(7)
                slide.shapes.add_picture(temp_img_path, left, top, height=height)
                
            finally:
                # 임시 파일 삭제 (빠른 처리)
                if temp_img_path and os.path.exists(temp_img_path):
                    try:
                        os.unlink(temp_img_path)
                    except (OSError, PermissionError) as e:
                        print(f"임시 파일 삭제 실패 (무시됨): {e}")
                        # 임시 파일 삭제 실패는 무시하고 계속 진행
        
        # PPTX 파일 저장
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"변환 중 오류 발생: {str(e)}")
        return False

# 파일 크기 초과 오류 처리
@app.errorhandler(413)
def too_large(e):
    flash('파일 크기가 100MB를 초과합니다. 더 작은 파일을 선택해주세요.')
    return redirect(url_for('index'))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        print("파일 업로드 요청 시작")
        
        # 1단계: 파일 존재 여부 확인
        if 'file' not in request.files:
            flash('파일이 선택되지 않았습니다.')
            return redirect(request.url)
        
        file = request.files['file']
        
        # 2단계: 파일명 확인
        if file.filename == '':
            flash('파일이 선택되지 않았습니다.')
            return redirect(request.url)
        
        # 3단계: 파일 크기 확인
        file.seek(0, 2)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > 100 * 1024 * 1024:  # 100MB
            flash(f'파일 크기가 너무 큽니다. (현재: {file_size // (1024*1024)}MB, 최대: 100MB)')
            return redirect(request.url)
        
        print(f"파일 크기: {file_size // (1024*1024)}MB")
        
        # 4단계: 파일 형식 확인 및 처리
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_ext = filename.rsplit('.', 1)[1].lower()
            input_path = os.path.join(UPLOAD_FOLDER, filename)
            
            print(f"파일 저장 중 - {input_path}")
            try:
                file.save(input_path)
                print("파일 저장 완료")
            except Exception as e:
                flash(f'파일 저장 중 오류가 발생했습니다: {str(e)}')
                return redirect(url_for('index'))
            
            # 변환 처리
            conversion_success = False
            output_path = None
            
            if file_ext == 'pdf':
                # PDF → DOCX 변환
                output_filename = filename.rsplit('.', 1)[0] + '.docx'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                quality = request.form.get('quality', 'medium')
                print(f"PDF → DOCX 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    conversion_success = pdf_to_docx(input_path, output_path, quality)
                except Exception as e:
                    print(f"변환 중 예외 발생: {str(e)}")
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
                    
            elif file_ext == 'docx':
                # DOCX → PDF 변환
                output_filename = filename.rsplit('.', 1)[0] + '.pdf'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                print(f"DOCX → PDF 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    conversion_success = docx_to_pdf(input_path, output_path)
                except Exception as e:
                    print(f"변환 중 예외 발생: {str(e)}")
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
            
            # 변환 결과 처리
            if conversion_success:
                print("변환 성공 - 다운로드 준비")
                
                # 업로드된 파일 정리
                try:
                    os.remove(input_path)
                    print("임시 파일 삭제 완료")
                except Exception as e:
                    print(f"임시 파일 삭제 실패 (무시됨): {e}")
                
                # 파일 다운로드 제공
                try:
                    print("파일 다운로드 시작")
                    return send_file(output_path, as_attachment=True, download_name=output_filename)
                except Exception as e:
                    print(f"파일 다운로드 오류: {str(e)}")
                    flash(f'파일 다운로드 중 오류가 발생했습니다: {str(e)}')
                    return redirect(url_for('index'))
            else:
                print("변환 실패 - 정리 작업")
                flash('파일 변환에 실패했습니다. 다시 시도해주세요.')
                
                # 실패한 파일들 정리
                for cleanup_path in [input_path, output_path]:
                    try:
                        if cleanup_path and os.path.exists(cleanup_path):
                            os.remove(cleanup_path)
                    except Exception as e:
                        print(f"파일 정리 실패 (무시됨): {e}")
                
                return redirect(url_for('index'))
        else:
            flash('PDF 또는 DOCX 파일만 업로드 가능합니다.')
            return redirect(url_for('index'))
            
    except Exception as e:
        print(f"업로드 처리 중 예외 발생: {str(e)}")
        flash('파일 처리 중 오류가 발생했습니다.')
        return redirect(url_for('index'))

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)